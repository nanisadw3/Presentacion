import customtkinter as ctk
import pandas as pd
from tkinter import filedialog, messagebox
import os
from CTkTable import CTkTable
import threading
import traceback
from termcolor import colored
import db_helper

class ExcelViewerApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("Visor de Excel (Nativo de CustomTkinter)")
        self.geometry("1500x850")
        
        # Configuración del tema
        ctk.set_appearance_mode("Dark")
        ctk.set_default_color_theme("blue")

        # Contenedor principal
        self.main_frame = ctk.CTkFrame(self, fg_color="transparent")
        self.main_frame.pack(fill="both", expand=True, padx=20, pady=20)

        # Header Frame
        self.top_frame = ctk.CTkFrame(self.main_frame, corner_radius=10)
        self.top_frame.pack(fill="x", pady=(0, 20))

        self.btn_buscar = ctk.CTkButton(self.top_frame, 
                                        text="Buscar Archivo Excel", 
                                        font=("Roboto", 14, "bold"),
                                        command=self.load_excel,
                                        height=40)
        self.btn_buscar.pack(pady=15, padx=20, side="left")

        self.lbl_file = ctk.CTkLabel(self.top_frame, 
                                     text="Ningún archivo seleccionado",
                                     font=("Roboto", 14))
        self.lbl_file.pack(pady=15, padx=20, side="left")

        # Botón para guardar las tablas en la Base de Datos SQLite
        self.btn_guardar = ctk.CTkButton(self.top_frame, 
                                        text="Guardar en Base de Datos", 
                                        font=("Roboto", 14, "bold"),
                                        command=self.save_to_database,
                                        height=40,
                                        fg_color="#28a745", hover_color="#218838")
        self.btn_guardar.pack(pady=15, padx=10, side="left")

        # Botón para mandar la información a PowerPoint
        self.btn_powerpoint = ctk.CTkButton(self.top_frame, 
                                            text="Mandar a PowerPoint", 
                                            font=("Roboto", 14, "bold"),
                                            command=self.send_to_powerpoint,
                                            height=40,
                                            fg_color="#8b5cf6", hover_color="#7c3aed")
        self.btn_powerpoint.pack(pady=15, padx=10, side="left")

        # Botón para agregar años manualmente a la BD
        self.btn_add_year = ctk.CTkButton(self.top_frame, 
                                          text="Agregar Año Extra", 
                                          font=("Roboto", 14, "bold"),
                                          command=self.open_add_year_dialog,
                                          height=40,
                                          fg_color="#007bff", hover_color="#0056b3")
        self.btn_add_year.pack(pady=15, padx=10, side="left")

        # ComboBox para alternar visualización de procesos
        self.lbl_proceso = ctk.CTkLabel(self.top_frame, text="Proceso:", font=("Roboto", 14, "bold"))
        self.lbl_proceso.pack(pady=15, padx=(20, 5), side="left")
        
        self.cb_proceso = ctk.CTkComboBox(self.top_frame, 
                                            values=["Crudo", "Gasolinas", "Diesel", "Turbosina", "Asfalto", "Combustoleo", "Cadereyta -Crudo", "Cadereyta -Gasolinas", "Cadereyta -Diesel"],
                                            font=("Roboto", 14),
                                            command=self.on_proceso_changed,
                                            state="readonly",
                                            width=150)
        self.cb_proceso.pack(pady=15, padx=5, side="left")
        self.cb_proceso.set("Crudo")


        # Scrollable Frame para contener la tabla
        self.scroll_frame = ctk.CTkScrollableFrame(self.main_frame, corner_radius=10)
        self.scroll_frame.pack(fill="both", expand=True)

        self.table = None
        self.df_data = None # Para la primera tabla
        self.df_snr = None  # Para la segunda tabla (AE-AF)
        self.df_prod = None # Para la tercera tabla (AE-AF, 21-40)
        self.df_sim = None  # Para la cuarta tabla (Simulación 12 meses)
        
        # Datos de Gasolinas
        self.df_data_gasolinas = None
        self.df_snr_gasolinas = None
        self.df_prod_gasolinas = None
        self.df_sim_gasolinas = None

        # Datos de Diesel
        self.df_data_diesel = None
        self.df_snr_diesel = None
        self.df_prod_diesel = None
        self.df_sim_diesel = None
 
        # Datos de Cadereyta -Gasolinas
        self.df_data_cad_gas = None
        self.df_snr_cad_gas = None
        self.df_prod_cad_gas = None
        self.df_sim_cad_gas = None

        # Datos de Cadereyta -Diesel
        self.df_data_cad_die = None
        self.df_snr_cad_die = None
        self.df_prod_cad_die = None
        self.df_sim_cad_die = None
        # Datos de Cadereyta -Crudo
        self.df_data_cad = None
        self.df_snr_cad = None
        self.df_prod_cad = None
        self.df_sim_cad = None
 
        # Datos de Turbosina
        self.df_data_turbosina = None

        self.df_snr_turbosina = None
        self.df_prod_turbosina = None
        self.df_sim_turbosina = None

        # Definir directorio por defecto (con fallback si no existe el de descargas)
        downloads_path = "/mnt/c/Users/10900096799/Downloads"
        if os.path.exists(downloads_path):
            self.default_dir = downloads_path
        else:
            self.default_dir = os.path.dirname(os.path.abspath(__file__))

        # Barra de progreso (inicialmente oculta)
        self.progress_bar = ctk.CTkProgressBar(self.top_frame, width=200)
        self.progress_bar.set(0.0)

    def set_loading_state(self, is_loading, loading_text=""):
        if is_loading:
            self.lbl_file.configure(text=loading_text)
            self.btn_buscar.configure(state="disabled")
            self.btn_guardar.configure(state="disabled")
            self.btn_powerpoint.configure(state="disabled")
            self.progress_bar.pack(pady=15, padx=20, side="left")
            self.progress_bar.set(0.0)
        else:
            self.btn_buscar.configure(state="normal")
            self.btn_guardar.configure(state="normal")
            self.btn_powerpoint.configure(state="normal")
            self.progress_bar.pack_forget()

    def update_progress(self, val, text_msg=None):
        self.progress_bar.set(val)
        if text_msg is not None:
            self.lbl_file.configure(text=text_msg)

    def open_add_year_dialog(self):
        dialog = ctk.CTkToplevel(self)
        dialog.title("Agregar/Editar Producción")
        dialog.geometry("450x500")
        dialog.transient(self)
        dialog.grab_set()

        lbl_proceso = ctk.CTkLabel(dialog, text="Proceso:", font=("Roboto", 14, "bold"))
        lbl_proceso.pack(pady=(10, 0))
        opciones_procesos = ["Crudo", "Cadereyta -Crudo", "Gasolinas", "Cadereyta -Gasolinas", "Diesel", "Cadereyta -Diesel", "Turbosina", "Asfalto", "Combustoleo"]
        combo_proceso = ctk.CTkComboBox(dialog, values=opciones_procesos, width=250)
        combo_proceso.pack(pady=(5, 10))

        tipo_var = ctk.StringVar(value="mes")
        frame_radios = ctk.CTkFrame(dialog, fg_color="transparent")
        frame_radios.pack(pady=10)
        
        radio_mes = ctk.CTkRadioButton(frame_radios, text="Agregar un Mes", variable=tipo_var, value="mes")
        radio_mes.pack(side="left", padx=10)
        radio_anio = ctk.CTkRadioButton(frame_radios, text="Agregar un Año (Total)", variable=tipo_var, value="anio")
        radio_anio.pack(side="left", padx=10)

        lbl_anio = ctk.CTkLabel(dialog, text="Año (ej. 2024):", font=("Roboto", 14, "bold"))
        lbl_anio.pack(pady=(10, 0))
        entry_anio = ctk.CTkEntry(dialog, width=250)
        entry_anio.pack(pady=(5, 10))

        frame_mes_input = ctk.CTkFrame(dialog, fg_color="transparent")
        frame_mes_input.pack(pady=5)
        lbl_mes = ctk.CTkLabel(frame_mes_input, text="Mes:", font=("Roboto", 14, "bold"))
        lbl_mes.pack(side="left", padx=5)
        meses_opciones = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
        combo_mes = ctk.CTkComboBox(frame_mes_input, values=meses_opciones, width=150)
        combo_mes.pack(side="left", padx=5)

        def on_radio_change(*args):
            if tipo_var.get() == "anio":
                combo_mes.configure(state="disabled")
            else:
                combo_mes.configure(state="normal")
        tipo_var.trace("w", on_radio_change)

        lbl_prod = ctk.CTkLabel(dialog, text="Producción:", font=("Roboto", 14, "bold"))
        lbl_prod.pack(pady=(10, 0))
        entry_prod = ctk.CTkEntry(dialog, width=250)
        entry_prod.pack(pady=(5, 20))

        def on_save():
            proceso = combo_proceso.get()
            anio = entry_anio.get().strip()
            if not anio.isdigit():
                messagebox.showerror("Error", "El año debe ser un número válido.")
                return
            
            try:
                prod = float(entry_prod.get().strip())
            except:
                messagebox.showerror("Error", "La producción debe ser numérica.")
                return
            
            import db_helper
            mes_val = combo_mes.get() if tipo_var.get() == "mes" else "AÑO"
            
            db_helper.save_extra_prod(proceso, anio, mes_val, prod)
            
            tipo_txt = "mes de " + mes_val if mes_val != "AÑO" else "año completo"
            messagebox.showinfo("Éxito", f"Producción del {tipo_txt} para '{proceso}' en {anio} guardada correctamente.\n\nVuelve a hacer clic en 'Buscar Archivo Excel' para recargar las gráficas.")
            dialog.destroy()

        def on_clear():
            if messagebox.askyesno("Confirmar", "¿Estás seguro de que quieres borrar TODOS los años y meses extra de la Base de Datos? (Esto no afecta al Excel original)"):
                import db_helper
                db_helper.clear_db()
                messagebox.showinfo("Limpieza Completa", "Se ha borrado toda la información extra de la base de datos local.\n\nVuelve a hacer clic en 'Buscar Archivo Excel' para recargar las gráficas.")
                dialog.destroy()

        btn_save = ctk.CTkButton(dialog, text="Guardar", command=on_save, fg_color="#28a745", hover_color="#218838")
        btn_save.pack(pady=(15, 5))

        btn_clear = ctk.CTkButton(dialog, text="Limpiar Base de Datos", command=on_clear, fg_color="#dc3545", hover_color="#c82333")
        btn_clear.pack(pady=5)

    def save_to_database(self):
        if self.df_data is None or self.df_snr is None or self.df_prod is None or self.df_sim is None:
            messagebox.showerror("Error", "Primero debes buscar un archivo Excel para extraer los datos.")
            return

        db_path = filedialog.asksaveasfilename(
            title="Crear o seleccionar Base de Datos SQLite",
            initialdir=self.default_dir,
            initialfile="datos_extraidos.db",
            defaultextension=".db",
            filetypes=[("SQLite Database", "*.db *.sqlite")]
        )

        if db_path:
            import sqlite3
            import math
            try:
                conn = sqlite3.connect(db_path)

                # Redondear todos los valores numéricos a enteros para la BD
                # (la pantalla conserva los decimales de MCP y Producción).
                def to_int_val(v):
                    try:
                        if v == "" or v is None:
                            return v
                        if isinstance(v, str):
                            return v
                        f = float(v)
                        if math.isinf(f) or math.isnan(f):
                            return ""
                        return int(round(f))
                    except (ValueError, TypeError):
                        return v

                def redondear_df(df):
                    out = df.copy()
                    
                    # Renombrar columnas vacías o duplicadas para que SQLite no arroje error
                    nuevas_columnas = []
                    for i, c in enumerate(out.columns):
                        if str(c).strip() == "" or pd.isna(c):
                            nuevas_columnas.append(f"Col_Vacia_{i+1}")
                        else:
                            nuevas_columnas.append(str(c))
                            
                    vistos = set()
                    for i in range(len(nuevas_columnas)):
                        orig = nuevas_columnas[i]
                        if orig in vistos:
                            j = 1
                            while f"{orig}_{j}" in vistos:
                                j += 1
                            nuevas_columnas[i] = f"{orig}_{j}"
                        vistos.add(nuevas_columnas[i])
                        
                    out.columns = nuevas_columnas

                    for c in out.columns:
                        out[c] = out[c].map(to_int_val)
                    return out

                def save_df(df, table_name):
                    if df is not None and len(df.columns) > 0:
                        redondear_df(df).to_sql(table_name, conn, if_exists='replace', index=False)

                # Guardar datos de Crudo
                save_df(self.df_data, 'crudo_tabla_principal')
                save_df(self.df_snr, 'crudo_programa_snr')
                save_df(self.df_prod, 'crudo_produccion')
                save_df(self.df_sim, 'crudo_simulacion_anual')

                # Guardar datos de Crudo Cadereyta
                save_df(getattr(self, 'df_data_cad', None), 'cadereyta_crudo_tabla_principal')
                save_df(getattr(self, 'df_snr_cad', None), 'cadereyta_crudo_programa_snr')
                save_df(getattr(self, 'df_prod_cad', None), 'cadereyta_crudo_produccion')
                save_df(getattr(self, 'df_sim_cad', None), 'cadereyta_crudo_simulacion_anual')

                 # Guardar datos de Gasolinas
                save_df(self.df_data_gasolinas, 'gasolinas_tabla_principal')
                save_df(self.df_snr_gasolinas, 'gasolinas_programa_snr')
                save_df(self.df_prod_gasolinas, 'gasolinas_produccion')
                save_df(getattr(self, 'df_sim_gasolinas', None), 'gasolinas_simulacion_anual')

                # Guardar datos de Diesel
                save_df(self.df_data_diesel, 'diesel_tabla_principal')
                save_df(self.df_snr_diesel, 'diesel_programa_snr')
                save_df(self.df_prod_diesel, 'diesel_produccion')
                save_df(getattr(self, 'df_sim_diesel', None), 'diesel_simulacion_anual')

                # Guardar datos de Turbosina
                save_df(self.df_data_turbosina, 'turbosina_tabla_principal')
                save_df(self.df_snr_turbosina, 'turbosina_programa_snr')
                save_df(self.df_prod_turbosina, 'turbosina_produccion')
                save_df(getattr(self, 'df_sim_turbosina', None), 'turbosina_simulacion_anual')

                # Guardar datos de Asfalto
                save_df(getattr(self, 'df_data_asfalto', None), 'asfalto_tabla_principal')
                save_df(getattr(self, 'df_snr_asfalto', None), 'asfalto_programa_snr')
                save_df(getattr(self, 'df_prod_asfalto', None), 'asfalto_produccion')
                save_df(getattr(self, 'df_sim_asfalto', None), 'asfalto_simulacion_anual')

                # Guardar datos de Combustoleo
                save_df(getattr(self, 'df_data_combustoleo', None), 'combustoleo_tabla_principal')
                save_df(getattr(self, 'df_snr_combustoleo', None), 'combustoleo_programa_snr')
                save_df(getattr(self, 'df_prod_combustoleo', None), 'combustoleo_produccion')
                save_df(getattr(self, 'df_sim_combustoleo', None), 'combustoleo_simulacion_anual')

                # Guardar datos de Cadereyta - Gasolinas
                save_df(getattr(self, 'df_data_cad_gas', None), 'cadereyta_gasolinas_tabla_principal')
                save_df(getattr(self, 'df_snr_cad_gas', None), 'cadereyta_gasolinas_programa_snr')
                save_df(getattr(self, 'df_prod_cad_gas', None), 'cadereyta_gasolinas_produccion')
                save_df(getattr(self, 'df_sim_cad_gas', None), 'cadereyta_gasolinas_simulacion_anual')

                # Guardar datos de Cadereyta - Diesel
                save_df(getattr(self, 'df_data_cad_die', None), 'cadereyta_diesel_tabla_principal')
                save_df(getattr(self, 'df_snr_cad_die', None), 'cadereyta_diesel_programa_snr')
                save_df(getattr(self, 'df_prod_cad_die', None), 'cadereyta_diesel_produccion')
                save_df(getattr(self, 'df_sim_cad_die', None), 'cadereyta_diesel_simulacion_anual')

                conn.commit()
                conn.close()
                messagebox.showinfo("Éxito", f"¡Los datos de todos los procesos han sido guardados en la base de datos!\n\nRuta:\n{db_path}")
            except Exception as e:
                messagebox.showerror("Error", f"Ocurrió un error al guardar en SQLite:\n{str(e)}")

    def load_excel(self):
        file_path = filedialog.askopenfilename(
            title="Seleccionar archivo Excel",
            initialdir=self.default_dir,
            filetypes=[("Archivos de Excel", "*.xlsx *.xls *.xlsm")]
        )

        if file_path:
            self.set_loading_state(True, "Cargando datos del Excel...")
            # Correr carga en segundo plano
            threading.Thread(target=self.async_load_data, args=(file_path,), daemon=True).start()

    def async_load_data(self, file_path):
        try:
            import warnings
            warnings.filterwarnings('ignore', category=UserWarning, module='openpyxl')

            # Función para eliminar filas donde la producción (última columna) sea 0 o vacía
            def filter_zero_rows(df):
                if df.empty: return df
                def is_not_zero(row):
                    try:
                        val = row.iloc[-1]
                        if val == "" or val is None: return False
                        return float(val) != 0
                    except:
                        return True
                return df[df.apply(is_not_zero, axis=1)]

            # Función segura para quitar decimales sin causar errores de tipo
            def remove_decimals(df_to_clean, skip_first=False, skip_last=False):
                def safe_round(val):
                    try:
                        if pd.isna(val):
                            return ""
                    except ValueError:
                        return ""
                    try:
                        f_val = float(val)
                        import math
                        if math.isinf(f_val) or math.isnan(f_val):
                            return ""
                        return int(round(f_val))
                    except:
                        return str(val).strip()

                def keep_raw(val):
                    try:
                        if pd.isna(val):
                            return ""
                    except ValueError:
                        return ""
                    try:
                        f_val = float(val)
                        import math
                        if math.isinf(f_val) or math.isnan(f_val):
                            return ""
                        return val
                    except:
                        return str(val).strip()

                if skip_first and df_to_clean.shape[1] > 1:
                    first_col = df_to_clean.iloc[:, 0].map(keep_raw)
                    rest = df_to_clean.iloc[:, 1:]
                    if hasattr(rest, 'map'):
                        rest = rest.map(safe_round)
                    else:
                        rest = rest.applymap(safe_round)
                    return pd.concat([first_col, rest], axis=1)

                if skip_last and df_to_clean.shape[1] > 1:
                    rest = df_to_clean.iloc[:, :-1]
                    last_col = df_to_clean.iloc[:, -1].map(keep_raw)
                    if hasattr(rest, 'map'):
                        rest = rest.map(safe_round)
                    else:
                        rest = rest.applymap(safe_round)
                    return pd.concat([rest, last_col], axis=1)

                if hasattr(df_to_clean, 'map'):
                    return df_to_clean.map(safe_round)
                else:
                    return df_to_clean.applymap(safe_round)

            # Autodetectar hoja
            xls = pd.ExcelFile(file_path)
            sheet_to_use = None
            for sheet in xls.sheet_names:
                normalizado = sheet.lower().replace(" ", "").replace("_", "")
                if "enviocalculopromedio" in normalizado:
                    sheet_to_use = sheet
                    break
            if not sheet_to_use:
                for sheet in xls.sheet_names:
                    normalizado = sheet.lower().replace(" ", "").replace("_", "")
                    if "enviocalculo" in normalizado:
                        sheet_to_use = sheet
                        break
            if not sheet_to_use:
                sheet_to_use = xls.sheet_names[0]

            # --- LEER LA HOJA COMPLETA UNA SOLA VEZ ---
            df_sheet = pd.read_excel(file_path, sheet_name=sheet_to_use, header=None)
            self.after(0, self.update_progress, 0.1, "Leyendo hoja de cálculo...")

            # Helper para formatear encabezados
            def get_clean_headers(row_idx, start_col, end_col):
                headers = []
                for val in df_sheet.iloc[row_idx, start_col:end_col]:
                    if pd.isna(val):
                        headers.append("")
                    else:
                        try:
                            headers.append(str(int(round(float(val)))))
                        except:
                            headers.append(str(val))
                return headers

            def merge_extra_prod(proceso_name, df_prod_current):
                extra_rows = db_helper.get_extra_prod(proceso_name)
                
                cleaned_data = []
                for _, row in df_prod_current.iterrows():
                    val0 = str(row.iloc[0]).strip()
                    if val0.endswith('.0'):
                        val0 = val0[:-2]
                    if val0.lower() in ['', 'nan', 'nat', 'none']:
                        continue
                    cleaned_data.append((val0, row.iloc[1]))

                if not extra_rows:
                    if not cleaned_data:
                        return df_prod_current
                    return pd.DataFrame(cleaned_data, columns=["Año/Mes", "Produccion"])

                for r in extra_rows:
                    anio, mes, prod = r
                    anio = str(anio).strip()
                    mes = str(mes).strip()
                    
                    if mes == "AÑO":
                        insert_pos = 0
                        replaced = False
                        for i, (c_name, c_val) in enumerate(cleaned_data):
                            if c_name.isdigit() and len(c_name) == 4:
                                if c_name == anio:
                                    cleaned_data[i] = (anio, prod)
                                    replaced = True
                                    break
                                elif int(c_name) > int(anio):
                                    insert_pos = i
                                    break
                                else:
                                    insert_pos = i + 1
                            else:
                                insert_pos = i
                                break
                                
                        if not replaced:
                            cleaned_data.insert(insert_pos, (anio, prod))
                    else:
                        month_block_start = len(cleaned_data)
                        for i, (c_name, c_val) in enumerate(cleaned_data):
                            if not (c_name.isdigit() and len(c_name) == 4):
                                month_block_start = i
                                break
                                
                        meses_orden = ["ene", "feb", "mar", "abr", "may", "jun", "jul", "ago", "sep", "oct", "nov", "dic"]
                        try:
                            target_month_idx = meses_orden.index(mes.lower()[:3])
                        except ValueError:
                            target_month_idx = 99
                            
                        insert_pos = month_block_start
                        replaced = False
                        
                        for i in range(month_block_start, len(cleaned_data)):
                            c_name = cleaned_data[i][0]
                            curr_month_str = c_name.lower()[:3]
                            try:
                                curr_m_idx = meses_orden.index(curr_month_str)
                            except ValueError:
                                curr_m_idx = -1
                                
                            if curr_m_idx == target_month_idx:
                                cleaned_data[i] = (mes, prod)
                                replaced = True
                                break
                            elif curr_m_idx > target_month_idx:
                                insert_pos = i
                                break
                            else:
                                insert_pos = i + 1
                                
                        if not replaced:
                            cleaned_data.insert(insert_pos, (mes, prod))
                            
                df_merged = pd.DataFrame(cleaned_data, columns=["Año/Mes", "Produccion"])
                return df_merged

            # --- 1. PROCESAR CRUDO (Cadereyta) ---
            self.after(0, self.update_progress, 0.2, "Procesando Crudo Cadereyta...")
            # --- 1. PROCESAR CRUDO (Estándar) ---
            self.after(0, self.update_progress, 0.2, "Procesando Crudo...")
            # Leer Tabla 1 (Rows 21-51 -> index 20:51), Cols A-H (0:8)
            headers_cad = get_clean_headers(0, 0, 8)
            df_data = df_sheet.iloc[20:51, 0:8].copy()
            df_data.columns = headers_cad
            df_data = df_data.dropna(how='all')
            df_data = remove_decimals(df_data)
            df_data = filter_zero_rows(df_data)

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Cols AE-AF (30:32)
            df_snr = df_sheet.iloc[73:104, 30:32].copy()
            df_snr.columns = ["CMP", "PODIM"]
            df_snr = df_snr.dropna(how='all').dropna(axis=1, how='all')
            df_snr = remove_decimals(df_snr, skip_first=True)
            df_snr_copy = df_snr.copy()

            # Leer Tabla 3 (Fecha y Producción, Rows 21-40 -> index 20:40), Cols AE-AF (30:32)
            df_prod_raw = df_sheet.iloc[20:40, 30:32].copy()
            df_prod_raw.columns = ["Año/Mes", "Produccion"]
            df_prod_raw = df_prod_raw.dropna(how='all')
            
            dic_idx_crudo = -1
            for idx, row in df_prod_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_crudo = idx - 20
                    break
            
            if dic_idx_crudo != -1:
                df_prod = df_prod_raw.iloc[:dic_idx_crudo + 1]
            else:
                df_prod = df_prod_raw.iloc[:20]
                
            df_prod = df_prod.dropna(axis=1, how='all')
            df_prod = remove_decimals(df_prod, skip_last=True)
            df_prod = merge_extra_prod("Crudo", df_prod)
            self.df_prod = df_prod.copy()
            df_prod_copy = df_prod.copy()

            num_dias_reales = len(df_data) if not df_data.empty else 31

            # --- 1.0 PROCESAR CRUDO CADEREYTA (Específico para Diapositiva 10) ---
            # Columnas A (0) y C (2), Filas 21-51
            df_data_cad = df_sheet.iloc[20:51, [0, 2]].copy()
            df_data_cad.columns = ["Crudo", "Cadereyta"]
            df_data_cad = df_data_cad.dropna(how='all')
            df_data_cad = remove_decimals(df_data_cad)
            df_data_cad = filter_zero_rows(df_data_cad)

            # Programa: Columna BI (60) repetida, Filas 74-104
            df_snr_cad = df_sheet.iloc[73:104, [60, 60]].copy()
            df_snr_cad.columns = ["CMP", "PODIM"]
            df_snr_cad = df_snr_cad.dropna(how='all').dropna(axis=1, how='all')
            df_snr_cad = remove_decimals(df_snr_cad, skip_first=True)
            df_snr_cad_copy = df_snr_cad.copy()
            
            # Fechas: Cols AW-AX (48:50), Filas 21-40
            df_prod_cad = df_sheet.iloc[20:40, 48:50].copy()
            df_prod_cad.columns = ["Año/Mes", "Produccion"]
            df_prod_cad = df_prod_cad.dropna(how='all')
            df_prod_cad = remove_decimals(df_prod_cad, skip_last=True)
            df_prod_cad = merge_extra_prod("Cadereyta -Crudo", df_prod_cad)
            df_prod_cad_copy = df_prod_cad.copy()
            
            # --- 1.1 PROCESAR GASOLINAS (Cadereyta) ---
            self.after(0, self.update_progress, 0.25, "Procesando Gasolinas Cadereyta...")
            headers_cad_gas = ["Cadereyta Gas - Día", "Cadereyta Gas - Producción"]
            # Leer Tabla 1 (Rows 22-51 -> index 21:51), Cols L y M (11, 12)
            df_gas_cad = df_sheet.iloc[21:51, [11, 12]].copy()
            df_gas_cad.columns = headers_cad_gas
            df_gas_cad = df_gas_cad.dropna(how='all')
            df_gas_cad = remove_decimals(df_gas_cad)
            df_gas_cad = filter_zero_rows(df_gas_cad)
            df_data_cad_gas = df_gas_cad.copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Col BS repetida (70, 70)
            df_snr_cad_gas = df_sheet.iloc[73:104, [70, 70]].copy()
            df_snr_cad_gas = df_snr_cad_gas.dropna(how='all').dropna(axis=1, how='all')
            df_snr_cad_gas = remove_decimals(df_snr_cad_gas, skip_first=True)
            df_snr_cad_gas_copy = df_snr_cad_gas.copy()
            # Aplicar recorte de días reales basado en Crudo
            df_data_cad_gas = df_data_cad_gas.iloc[:num_dias_reales]

            # Leer Tabla 3 (Fecha y Producción Cadereyta Gas, Rows 21-40 -> index 20:40), Cols BO:BP (66:68)
            df_prod_cad_gas_raw = df_sheet.iloc[20:40, 66:68].copy()
            df_prod_cad_gas_raw = df_prod_cad_gas_raw.dropna(how='all')
            
            dic_idx_cad_gas = -1
            for idx, row in df_prod_cad_gas_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_cad_gas = idx - 20 
                    break
            
            if dic_idx_cad_gas != -1:
                df_prod_cad_gas = df_prod_cad_gas_raw.iloc[:dic_idx_cad_gas + 1]
            else:
                df_prod_cad_gas = df_prod_cad_gas_raw.iloc[:20]
                
            df_prod_cad_gas = df_prod_cad_gas.dropna(axis=1, how='all')
            df_prod_cad_gas = remove_decimals(df_prod_cad_gas, skip_last=True)
            df_prod_cad_gas = merge_extra_prod("Cadereyta -Gasolinas", df_prod_cad_gas)
            df_prod_cad_gas_copy = df_prod_cad_gas.copy()

            # --- 1.2 PROCESAR DIESEL (Cadereyta) ---
            self.after(0, self.update_progress, 0.3, "Procesando Diesel Cadereyta...")
            headers_cad_die = ["Cadereyta Die - Día", "Cadereyta Die - Producción"]
            # Leer Tabla 1 (Rows 74-104 -> index 73:104), Cols A y B (0, 1)
            df_die_cad = df_sheet.iloc[73:104, [0, 1]].copy()
            df_die_cad.columns = headers_cad_die
            df_die_cad = df_die_cad.dropna(how='all')
            df_die_cad = remove_decimals(df_die_cad)
            df_die_cad = filter_zero_rows(df_die_cad)
            df_data_cad_die = df_die_cad.copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Col CC repetida (80, 80)
            df_snr_cad_die = df_sheet.iloc[73:104, [80, 80]].copy()
            df_snr_cad_die = df_snr_cad_die.dropna(how='all').dropna(axis=1, how='all')
            df_snr_cad_die = remove_decimals(df_snr_cad_die, skip_first=True)
            df_snr_cad_die_copy = df_snr_cad_die.copy()
            # Aplicar recorte de días reales basado en Crudo
            df_data_cad_die = df_data_cad_die.iloc[:num_dias_reales]

            # Leer Tabla 3 (Fecha y Producción Cadereyta Die, Rows 21-40 -> index 20:40), Cols CG:CH (84:86)
            df_prod_cad_die_raw = df_sheet.iloc[20:40, 84:86].copy()
            df_prod_cad_die_raw = df_prod_cad_die_raw.dropna(how='all')
            
            dic_idx_cad_die = -1
            for idx, row in df_prod_cad_die_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_cad_die = idx - 20 
                    break
            
            if dic_idx_cad_die != -1:
                df_prod_cad_die = df_prod_cad_die_raw.iloc[:dic_idx_cad_die + 1]
            else:
                df_prod_cad_die = df_prod_cad_die_raw.iloc[:20]
                
            df_prod_cad_die = df_prod_cad_die.dropna(axis=1, how='all')
            df_prod_cad_die = remove_decimals(df_prod_cad_die, skip_last=True)
            df_prod_cad_die = merge_extra_prod("Cadereyta -Diesel", df_prod_cad_die)
            df_prod_cad_die_copy = df_prod_cad_die.copy()



            # --- 2. PROCESAR GASOLINAS ---
            self.after(0, self.update_progress, 0.35, "Procesando Gasolinas...")
            # Leer fila 1 (index 0) para encabezados, Cols L:S (11:19)
            clean_headers_gas = get_clean_headers(0, 11, 19)
            # Leer Tabla 1 (Rows 21-51 -> index 20:51), Cols L:S (11:19)
            df_gas = df_sheet.iloc[20:51, 11:19].copy()
            df_gas.columns = clean_headers_gas
            df_gas = df_gas.dropna(how='all').dropna(axis=1, how='all')
            df_gas = remove_decimals(df_gas)
            df_data_gasolinas = filter_zero_rows(df_gas).copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Cols AK:AL (36:38)
            df_snr_gas = df_sheet.iloc[73:104, 36:38].copy()
            df_snr_gas = df_snr_gas.dropna(how='all').dropna(axis=1, how='all')
            df_snr_gas = remove_decimals(df_snr_gas, skip_first=True)
            df_snr_gas_copy = df_snr_gas.copy()

            # Recortar filas de fin de mes
            df_data_gasolinas = df_data_gasolinas.iloc[:num_dias_reales]

            # Leer Tabla 3 (Años, Rows 21-120 -> index 20:120), Cols AG:AH (32:34)
            df_prod_gas_raw = df_sheet.iloc[20:120, 32:34].copy()
            df_prod_gas_raw = df_prod_gas_raw.dropna(how='all')
            
            dic_idx_gas = -1
            for idx, row in df_prod_gas_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_gas = idx - 20
                    break
            
            if dic_idx_gas != -1:
                df_prod_gas = df_prod_gas_raw.iloc[:dic_idx_gas + 1]
            else:
                df_prod_gas = df_prod_gas_raw.iloc[:20]
                
            df_prod_gas = df_prod_gas.dropna(axis=1, how='all')
            df_prod_gas = remove_decimals(df_prod_gas, skip_last=True)
            df_prod_gas = merge_extra_prod("Gasolinas", df_prod_gas)
            df_prod_gasolinas_copy = df_prod_gas.copy()


            # --- 3. PROCESAR DIESEL ---
            self.after(0, self.update_progress, 0.5, "Procesando Diesel...")
            # Leer fila 54 (index 53) para encabezados, Cols A:H (0:8)
            clean_headers_die = get_clean_headers(53, 0, 8)
            # Leer Tabla 1 (Rows 74-104 -> index 73:104), Cols A:H (0:8)
            df_die = df_sheet.iloc[73:104, 0:8].copy()
            df_die.columns = clean_headers_die
            df_die = df_die.dropna(how='all').dropna(axis=1, how='all')
            df_die = remove_decimals(df_die)
            df_data_diesel = filter_zero_rows(df_die).copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Cols AQ:AR (42:44)
            df_snr_die = df_sheet.iloc[73:104, 42:44].copy()
            df_snr_die = df_snr_die.dropna(how='all').dropna(axis=1, how='all')
            df_snr_die = remove_decimals(df_snr_die, skip_first=True)
            df_snr_die_copy = df_snr_die.copy()

            # Recortar filas de fin de mes
            df_data_diesel = df_data_diesel.iloc[:num_dias_reales]

            # Leer Tabla 3 (Años, Rows 21-120 -> index 20:120), Cols AK:AL (36:38)
            df_prod_die_raw = df_sheet.iloc[20:120, 36:38].copy()
            df_prod_die_raw = df_prod_die_raw.dropna(how='all')
            
            dic_idx_die = -1
            for idx, row in df_prod_die_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_die = idx - 20
                    break
            
            if dic_idx_die != -1:
                df_prod_die = df_prod_die_raw.iloc[:dic_idx_die + 1]
            else:
                df_prod_die = df_prod_die_raw.iloc[:20]
                
            df_prod_die = df_prod_die.dropna(axis=1, how='all')
            df_prod_die = remove_decimals(df_prod_die, skip_last=True)
            df_prod_die = merge_extra_prod("Diesel", df_prod_die)
            df_prod_diesel_copy = df_prod_die.copy()


            # --- 4. PROCESAR TURBOSINA ---
            self.after(0, self.update_progress, 0.62, "Procesando Turbosina...")
            # Leer fila 54 (index 53) para encabezados, Cols L:Q (11:17)
            clean_headers_turb = get_clean_headers(53, 11, 17)
            # Leer Tabla 1 (Rows 74-104 -> index 73:104), Cols L:Q (11:17)
            df_turb = df_sheet.iloc[73:104, 11:17].copy()
            df_turb.columns = clean_headers_turb
            df_turb = df_turb.dropna(how='all').dropna(axis=1, how='all')
            df_turb = remove_decimals(df_turb)
            df_data_turbosina = filter_zero_rows(df_turb).copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Cols AW:AX (48:50)
            df_snr_turb = df_sheet.iloc[73:104, 48:50].copy()
            df_snr_turb = df_snr_turb.dropna(how='all').dropna(axis=1, how='all')
            df_snr_turb = remove_decimals(df_snr_turb, skip_first=True)
            df_snr_turb_copy = df_snr_turb.copy()

            # Recortar filas de fin de mes
            df_data_turbosina = df_data_turbosina.iloc[:num_dias_reales]

            # Leer Tabla 3 (Años, Rows 21-120 -> index 20:120), Cols AM:AN (38:40)
            df_prod_turb_raw = df_sheet.iloc[20:120, 38:40].copy()
            df_prod_turb_raw = df_prod_turb_raw.dropna(how='all')
            
            dic_idx_turb = -1
            for idx, row in df_prod_turb_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_turb = idx - 20
                    break
            
            if dic_idx_turb != -1:
                df_prod_turb = df_prod_turb_raw.iloc[:dic_idx_turb + 1]
            else:
                df_prod_turb = df_prod_turb_raw.iloc[:20]
                
            df_prod_turb = df_prod_turb.dropna(axis=1, how='all')
            df_prod_turb = remove_decimals(df_prod_turb, skip_last=True)
            df_prod_turb = merge_extra_prod("Turbosina", df_prod_turb)
            df_prod_turbosina_copy = df_prod_turb.copy()


            # --- 5. PROCESAR ASFALTO ---
            self.after(0, self.update_progress, 0.74, "Procesando Asfalto...")
            # Leer Tabla 1 (Rows 122-152 -> index 121:152), Cols AM:AN (38:40)
            df_asf = df_sheet.iloc[121:152, 38:40].copy()
            df_asf.columns = ["Asfalto", "real"]
            df_asf = df_asf.dropna(how='all').dropna(axis=1, how='all')
            df_asf = remove_decimals(df_asf)
            df_data_asfalto = filter_zero_rows(df_asf).copy()

            # Leer Tabla 2 (Programa, Rows 122-152 -> index 121:152), Cols AK:AL (36:38)
            df_snr_asf = df_sheet.iloc[121:152, 36:38].copy()
            df_snr_asf = df_snr_asf.dropna(how='all').dropna(axis=1, how='all')
            df_snr_asf = remove_decimals(df_snr_asf, skip_first=True)
            df_snr_asf_copy = df_snr_asf.copy()

            # Recortar filas de fin de mes
            df_data_asfalto = df_data_asfalto.iloc[:num_dias_reales]

            # Leer Tabla 3 (Años, Rows 121-140 -> index 120:140), Cols AQ:AR (42:44)
            df_prod_asf_raw = df_sheet.iloc[120:140, 42:44].copy()
            df_prod_asf_raw = df_prod_asf_raw.dropna(how='all')

            dic_idx_asf = -1
            for idx, row in df_prod_asf_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_asf = idx - 120
                    break

            if dic_idx_asf != -1:
                df_prod_asf = df_prod_asf_raw.iloc[:dic_idx_asf + 1]
            else:
                df_prod_asf = df_prod_asf_raw.iloc[:20]

            df_prod_asf = df_prod_asf.dropna(axis=1, how='all')
            df_prod_asf = remove_decimals(df_prod_asf, skip_last=True)
            df_prod_asf = merge_extra_prod("Asfalto", df_prod_asf)
            df_prod_asfalto_copy = df_prod_asf.copy()


            # --- 6. PROCESAR COMBUSTOLEO ---
            self.after(0, self.update_progress, 0.85, "Procesando Combustoleo...")
            # Leer Tabla 1 (Rows 74-104 -> index 73:104), Cols U:V (20:22)
            df_comb = df_sheet.iloc[73:104, 20:22].copy()
            df_comb.columns = ["SNR", "Combustoleo"]
            df_comb = df_comb.dropna(how='all').dropna(axis=1, how='all')
            df_comb = remove_decimals(df_comb)
            df_data_combustoleo = filter_zero_rows(df_comb).copy()

            # Leer Tabla 2 (Programa, Rows 74-104 -> index 73:104), Cols BC:BD (54:56)
            df_snr_comb = df_sheet.iloc[73:104, 54:56].copy()
            df_snr_comb = df_snr_comb.dropna(how='all').dropna(axis=1, how='all')
            df_snr_comb = remove_decimals(df_snr_comb, skip_first=True)
            df_snr_comb_copy = df_snr_comb.copy()

            # Recortar filas de fin de mes
            df_data_combustoleo = df_data_combustoleo.iloc[:num_dias_reales]

            # Leer Tabla 3 (Años, Rows 21-40 -> index 20:40), Cols AQ:AR (42:44)
            df_prod_comb_raw = df_sheet.iloc[20:40, 42:44].copy()
            df_prod_comb_raw = df_prod_comb_raw.dropna(how='all')

            dic_idx_comb = -1
            for idx, row in df_prod_comb_raw.iterrows():
                val = str(row.iloc[0]).strip().lower()
                if "dic" in val or "diciembre" in val:
                    dic_idx_comb = idx - 20
                    break

            if dic_idx_comb != -1:
                df_prod_comb = df_prod_comb_raw.iloc[:dic_idx_comb + 1]
            else:
                df_prod_comb = df_prod_comb_raw.iloc[:20]

            df_prod_comb = df_prod_comb.dropna(axis=1, how='all')
            df_prod_comb = remove_decimals(df_prod_comb, skip_last=True)
            df_prod_comb = merge_extra_prod("Combustoleo", df_prod_comb)
            df_prod_combustoleo_copy = df_prod_comb.copy()


            # Crear Tabla 4 (Simulación 12 meses)
            import calendar
            from datetime import datetime
            now = datetime.now()
            current_year = now.year
            days_passed = now.timetuple().tm_yday

            meses_nombres = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
            meses_cortos = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
            dias_por_mes = [calendar.monthrange(current_year, m)[1] for m in range(1, 13)]

            # 1. Simulación Crudo
            prod_dict = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0
                
                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict[meses_nombres[i]] += p
                        break

            sim_data = []
            suma_total = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict[mes]
                total = prod * dias
                suma_total += total
                sim_data.append([mes, int(prod), dias, int(total)])
                
            promedio = suma_total / days_passed if days_passed > 0 else 0
            sim_data.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total)} | Prom: {promedio:.2f}"])
            df_sim = pd.DataFrame(sim_data, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])

            # 2. Simulación Gasolinas
            prod_dict_gas = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_gas.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0
                
                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_gas[meses_nombres[i]] += p
                        break

            sim_data_gas = []
            suma_total_gas = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_gas[mes]
                total = prod * dias
                suma_total_gas += total
                sim_data_gas.append([mes, int(prod), dias, int(total)])
                
            promedio_gas = suma_total_gas / days_passed if days_passed > 0 else 0
            sim_data_gas.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_gas)} | Prom: {promedio_gas:.2f}"])
            df_sim_gasolinas = pd.DataFrame(sim_data_gas, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])

            # 3. Simulación Diesel
            prod_dict_die = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_die.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0
                
                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_die[meses_nombres[i]] += p
                        break

            sim_data_die = []
            suma_total_die = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_die[mes]
                total = prod * dias
                suma_total_die += total
                sim_data_die.append([mes, int(prod), dias, int(total)])
                
            promedio_die = suma_total_die / days_passed if days_passed > 0 else 0
            sim_data_die.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_die)} | Prom: {promedio_die:.2f}"])
            df_sim_diesel = pd.DataFrame(sim_data_die, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])

            # 4. Simulación Turbosina
            prod_dict_turb = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_turb.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0
                
                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_turb[meses_nombres[i]] += p
                        break

            sim_data_turb = []
            suma_total_turb = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_turb[mes]
                total = prod * dias
                suma_total_turb += total
                sim_data_turb.append([mes, int(prod), dias, int(total)])
                
            promedio_turb = suma_total_turb / days_passed if days_passed > 0 else 0
            sim_data_turb.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_turb)} | Prom: {promedio_turb:.2f}"])
            df_sim_turbosina = pd.DataFrame(sim_data_turb, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])

            # 5. Simulación Asfalto
            prod_dict_asf = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_asf.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0

                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_asf[meses_nombres[i]] += p
                        break

            sim_data_asf = []
            suma_total_asf = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_asf[mes]
                total = prod * dias
                suma_total_asf += total
                sim_data_asf.append([mes, int(prod), dias, int(total)])

            promedio_asf = suma_total_asf / days_passed if days_passed > 0 else 0
            sim_data_asf.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_asf)} | Prom: {promedio_asf:.2f}"])
            df_sim_asfalto = pd.DataFrame(sim_data_asf, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])

            # 6. Simulación Combustoleo
            prod_dict_comb = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_comb.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try:
                    p = float(val_prod)
                except:
                    p = 0.0

                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_comb[meses_nombres[i]] += p
                        break

            sim_data_comb = []
            suma_total_comb = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_comb[mes]
                total = prod * dias
                suma_total_comb += total
                sim_data_comb.append([mes, int(prod), dias, int(total)])

            promedio_comb = suma_total_comb / days_passed if days_passed > 0 else 0
            sim_data_comb.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_comb)} | Prom: {promedio_comb:.2f}"])
            df_sim_combustoleo = pd.DataFrame(sim_data_comb, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])
 
            # 7. Simulación Cadereyta Gasolinas
            prod_dict_cad_gas = {m: 0.0 for m in meses_nombres}
            if df_prod_cad_gas is not None:
                for idx, row_data in df_prod_cad_gas.iterrows():
                    val_anio = str(row_data.iloc[0]).strip().lower()
                    val_prod = row_data.iloc[1]
                    try:
                        p = float(val_prod)
                    except:
                        p = 0.0
                    
                    for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                        if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                            prod_dict_cad_gas[meses_nombres[i]] += p
                            break
 
            sim_data_cad_gas = []
            suma_total_cad_gas = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_cad_gas[mes]
                total = prod * dias
                suma_total_cad_gas += total
                sim_data_cad_gas.append([mes, int(prod), dias, int(total)])
                
            promedio_cad_gas = suma_total_cad_gas / days_passed if days_passed > 0 else 0
            sim_data_cad_gas.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_cad_gas)} | Prom: {promedio_cad_gas:.2f}"])
            df_sim_cad_gas = pd.DataFrame(sim_data_cad_gas, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])
 
            # 8. Simulación Cadereyta Diesel
            prod_dict_cad_die = {m: 0.0 for m in meses_nombres}
            if df_prod_cad_die is not None:
                for idx, row_data in df_prod_cad_die.iterrows():
                    val_anio = str(row_data.iloc[0]).strip().lower()
                    val_prod = row_data.iloc[1]
                    try:
                        p = float(val_prod)
                    except:
                        p = 0.0
                    
                    for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                        if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                            prod_dict_cad_die[meses_nombres[i]] += p
                            break
 
            sim_data_cad_die = []
            suma_total_cad_die = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_cad_die[mes]
                total = prod * dias
                suma_total_cad_die += total
                sim_data_cad_die.append([mes, int(prod), dias, int(total)])
                
            promedio_cad_die = suma_total_cad_die / days_passed if days_passed > 0 else 0
            sim_data_cad_die.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_cad_die)} | Prom: {promedio_cad_die:.2f}"])
            df_sim_cad_die = pd.DataFrame(sim_data_cad_die, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])
            # Simulación Crudo Cadereyta
            prod_dict_cad = {m: 0.0 for m in meses_nombres}
            for idx, row_data in df_prod_cad.iterrows():
                val_anio = str(row_data.iloc[0]).strip().lower()
                val_prod = row_data.iloc[1]
                try: p = float(val_prod)
                except: p = 0.0
                for i, (m_largo, m_corto) in enumerate(zip(meses_nombres, meses_cortos)):
                    if m_largo.lower() in val_anio or m_corto.lower() in val_anio:
                        prod_dict_cad[meses_nombres[i]] += p
                        break
            sim_data_cad = []
            suma_total_cad = 0.0
            for i, mes in enumerate(meses_nombres):
                dias = dias_por_mes[i]
                prod = prod_dict_cad[mes]
                total = prod * dias
                suma_total_cad += total
                sim_data_cad.append([mes, int(prod), dias, int(total)])
            promedio_cad = suma_total_cad / days_passed if days_passed > 0 else 0
            sim_data_cad.append(["TOTALES", "---", f"Días pasados: {days_passed}", f"Suma: {int(suma_total_cad)} | Prom: {promedio_cad:.2f}"])
            df_sim_cad = pd.DataFrame(sim_data_cad, columns=["Mes", "Producción", "Días", "Total (Prod x Días)"])
 
            self.after(0, self.update_progress, 0.97, "Finalizando...")

            # Pasar datos a la interfaz (main thread)
            self.after(0, self.on_load_success, file_path, df_data, df_snr_copy, df_prod_copy, df_sim,
                       df_data_gasolinas, df_snr_gas_copy, df_prod_gasolinas_copy, df_sim_gasolinas,
                       df_data_diesel, df_snr_die_copy, df_prod_diesel_copy, df_sim_diesel,
                       df_data_turbosina, df_snr_turb_copy, df_prod_turbosina_copy, df_sim_turbosina,
                       df_data_asfalto, df_snr_asf_copy, df_prod_asfalto_copy, df_sim_asfalto,
                       df_data_combustoleo, df_snr_comb_copy, df_prod_combustoleo_copy, df_sim_combustoleo,
                       df_data_cad_gas, df_snr_cad_gas_copy, df_prod_cad_gas_copy, df_sim_cad_gas,
                       df_data_cad_die, df_snr_cad_die_copy, df_prod_cad_die_copy, df_sim_cad_die,
                       df_data_cad, df_snr_cad_copy, df_prod_cad_copy, df_sim_cad)

        except Exception as e:
            err_details = traceback.format_exc()
            self.after(0, self.on_load_error, str(e), err_details)

    def on_load_success(self, file_path, df_data, df_snr, df_prod, df_sim,
                            df_data_gasolinas=None, df_snr_gasolinas=None, df_prod_gasolinas=None, df_sim_gasolinas=None,
                            df_data_diesel=None, df_snr_diesel=None, df_prod_diesel=None, df_sim_diesel=None,
                            df_data_turbosina=None, df_snr_turbosina=None, df_prod_turbosina=None, df_sim_turbosina=None,
                            df_data_asfalto=None, df_snr_asfalto=None, df_prod_asfalto=None, df_sim_asfalto=None,
                            df_data_combustoleo=None, df_snr_combustoleo=None, df_prod_combustoleo=None, df_sim_combustoleo=None,
                            df_data_cad_gas=None, df_snr_cad_gas=None, df_prod_cad_gas=None, df_sim_cad_gas=None,
                            df_data_cad_die=None, df_snr_cad_die=None, df_prod_cad_die=None, df_sim_cad_die=None,
                            df_data_cad=None, df_snr_cad=None, df_prod_cad=None, df_sim_cad=None):
        self.df_data = df_data

        self.df_snr = df_snr
        self.df_prod = df_prod
        self.df_sim = df_sim
        
        self.df_data_gasolinas = df_data_gasolinas
        self.df_snr_gasolinas = df_snr_gasolinas
        self.df_prod_gasolinas = df_prod_gasolinas
        self.df_sim_gasolinas = df_sim_gasolinas

        self.df_data_diesel = df_data_diesel
        self.df_snr_diesel = df_snr_diesel
        self.df_prod_diesel = df_prod_diesel
        self.df_sim_diesel = df_sim_diesel
 
        self.df_data_cad_gas = df_data_cad_gas
        self.df_snr_cad_gas = df_snr_cad_gas
        self.df_prod_cad_gas = df_prod_cad_gas
        self.df_sim_cad_gas = df_sim_cad_gas

        self.df_data_cad_die = df_data_cad_die
        self.df_snr_cad_die = df_snr_cad_die
        self.df_prod_cad_die = df_prod_cad_die
        self.df_sim_cad_die = df_sim_cad_die
 
        self.df_data_cad = df_data_cad
        self.df_snr_cad = df_snr_cad
        self.df_prod_cad = df_prod_cad
        self.df_sim_cad = df_sim_cad
 
        self.df_data_turbosina = df_data_turbosina

        self.df_snr_turbosina = df_snr_turbosina
        self.df_prod_turbosina = df_prod_turbosina
        self.df_sim_turbosina = df_sim_turbosina

        self.df_data_asfalto = df_data_asfalto
        self.df_snr_asfalto = df_snr_asfalto
        self.df_prod_asfalto = df_prod_asfalto
        self.df_sim_asfalto = df_sim_asfalto

        self.df_data_combustoleo = df_data_combustoleo
        self.df_snr_combustoleo = df_snr_combustoleo
        self.df_prod_combustoleo = df_prod_combustoleo
        self.df_sim_combustoleo = df_sim_combustoleo

        # Mostrar las tablas correspondientes a la selección actual del ComboBox
        self.show_proceso_tables(self.cb_proceso.get())

        self.set_loading_state(False)
        self.lbl_file.configure(text=f"Archivo: {os.path.basename(file_path)}")

    def on_proceso_changed(self, selection):
        if self.df_data is None:
            return
        self.show_proceso_tables(selection)

    def show_proceso_tables(self, selection):
        # Limpiar tablas actuales
        if self.table is not None:
            self.table.destroy()
            self.table = None
        if hasattr(self, 'lbl_table1') and self.lbl_table1 is not None:
            self.lbl_table1.destroy()
        if hasattr(self, 'table2') and self.table2 is not None:
            self.table2.destroy()
            self.table2 = None
        if hasattr(self, 'lbl_table2') and self.lbl_table2 is not None:
            self.lbl_table2.destroy()
        if hasattr(self, 'table3') and self.table3 is not None:
            self.table3.destroy()
            self.table3 = None
        if hasattr(self, 'lbl_table3') and self.lbl_table3 is not None:
            self.lbl_table3.destroy()
        if hasattr(self, 'table4') and self.table4 is not None:
            self.table4.destroy()
            self.table4 = None
        if hasattr(self, 'lbl_table4') and self.lbl_table4 is not None:
            self.lbl_table4.destroy()

        if selection == "Crudo":
            df_data = self.df_data
            df_snr = self.df_snr
            df_prod = self.df_prod
            df_sim = self.df_sim
            lbl2_txt = "Programa (CMP, PODIM)"
            lbl3_txt = "Fecha y Producción (Año/Mes, Producción)"
        elif selection == "Cadereyta -Crudo":
            df_data = self.df_data_cad
            df_snr = self.df_snr_cad
            df_prod = self.df_prod_cad
            df_sim = self.df_sim_cad
            lbl2_txt = "Programa Crudo Cadereyta (Col BI x2, Filas 74-104)"
            lbl3_txt = "Fecha y Producción Crudo Cadereyta (AW-AX, Filas 21-40)"
        elif selection == "Cadereyta -Gasolinas":
            df_data = self.df_data_cad_gas
            df_snr = self.df_snr_cad_gas
            df_prod = self.df_prod_cad_gas
            df_sim = self.df_sim_cad_gas
            lbl2_txt = "Programa de Gasolinas (Col BS x2, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Gasolinas (BO-BP, Filas 21-40)"
        elif selection == "Cadereyta -Diesel":
            df_data = self.df_data_cad_die
            df_snr = self.df_snr_cad_die
            df_prod = self.df_prod_cad_die
            df_sim = self.df_sim_cad_die
            lbl2_txt = "Programa de Diesel (Col CC x2, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Diesel (CG-CH, Filas 21-40)"
        elif selection == "Gasolinas":
            df_data = self.df_data_gasolinas
            df_snr = self.df_snr_gasolinas
            df_prod = self.df_prod_gasolinas
            df_sim = self.df_sim_gasolinas
            lbl2_txt = "Programa de Gasolinas (AK-AL, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Gasolinas (AG-AH, Filas 21-40)"
        elif selection == "Diesel":
            df_data = self.df_data_diesel
            df_snr = self.df_snr_diesel
            df_prod = self.df_prod_diesel
            df_sim = self.df_sim_diesel
            lbl2_txt = "Programa de Diesel (AQ-AR, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Diesel (AK-AL, Filas 21-40)"
        elif selection == "Turbosina":
            df_data = self.df_data_turbosina
            df_snr = self.df_snr_turbosina
            df_prod = self.df_prod_turbosina
            df_sim = self.df_sim_turbosina
            lbl2_txt = "Programa de Turbosina (AW-AX, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Turbosina (AM-AN, Filas 21-40)"
        elif selection == "Asfalto":
            df_data = self.df_data_asfalto
            df_snr = self.df_snr_asfalto
            df_prod = self.df_prod_asfalto
            df_sim = self.df_sim_asfalto
            lbl2_txt = "Programa de Asfalto (AK-AL, Filas 122-152)"
            lbl3_txt = "Fecha y Producción de Asfalto (AQ-AR, Filas 121-140)"
        elif selection == "Combustoleo":
            df_data = self.df_data_combustoleo
            df_snr = self.df_snr_combustoleo
            df_prod = self.df_prod_combustoleo
            df_sim = self.df_sim_combustoleo
            lbl2_txt = "Programa de Combustoleo (BC-BD, Filas 74-104)"
            lbl3_txt = "Fecha y Producción de Combustoleo (AQ-AR, Filas 21-40)"

        if df_data is None or df_snr is None or df_prod is None:
            return

        # Mostrar valores redondeados a enteros en la vista previa.
        # (los cálculos de simulación siguen usando los datos originales sin redondear)
        def _red(val):
            if val == "" or val is None:
                return val
            try:
                f = float(val)
                import math
                if math.isinf(f) or math.isnan(f):
                    return ""
                return int(round(f))
            except (ValueError, TypeError):
                return val

        def _red_df(d):
            out = d.copy()
            for c in out.columns:
                out[c] = out[c].map(_red)
            return out

        df_data_v = _red_df(df_data)
        df_snr_v = _red_df(df_snr)
        df_prod_v = _red_df(df_prod)

        # Dibujar Tabla 1
        self.lbl_table1 = ctk.CTkLabel(self.scroll_frame, text=selection, font=("Roboto", 16, "bold"), text_color="#3484F0")
        self.lbl_table1.pack(pady=(20, 5))
        headers = list(df_data_v.columns)
        rows = df_data_v.to_numpy().tolist()
        if len(rows) > 200: rows = rows[:200]
        table_values = [headers] + rows
        self.table = CTkTable(
            master=self.scroll_frame, 
            row=len(table_values), 
            column=len(table_values[0]), 
            values=table_values,
            header_color="#1f538d",
            colors=["#2a2a2a", "#242424"],
            hover_color="#3a3a3a",
        )
        self.table.pack(expand=True, fill="both", padx=10, pady=10)

        # Dibujar Tabla 2
        self.lbl_table2 = ctk.CTkLabel(self.scroll_frame, text=f"Programa de {selection}", font=("Roboto", 16, "bold"), text_color="#3484F0")
        self.lbl_table2.pack(pady=(20, 5))
        headers2 = ["MCP", "PODIM"]
        rows2 = df_snr_v.to_numpy().tolist()
        if len(rows2) > 200: rows2 = rows2[:200]
        table_values2 = [headers2] + rows2
        self.table2 = CTkTable(
            master=self.scroll_frame, 
            row=len(table_values2), 
            column=len(table_values2[0]), 
            values=table_values2,
            header_color="#1f538d",
            colors=["#2a2a2a", "#242424"],
            hover_color="#3a3a3a",
        )
        self.table2.pack(expand=True, fill="both", padx=10, pady=10)

        # Dibujar Tabla 3
        self.lbl_table3 = ctk.CTkLabel(self.scroll_frame, text=lbl3_txt, font=("Roboto", 16, "bold"), text_color="#3484F0")
        self.lbl_table3.pack(pady=(20, 5))
        headers3 = ["Año / Mes", "Producción"]
        rows3 = df_prod_v.to_numpy().tolist()
        if len(rows3) > 200: rows3 = rows3[:200]
        table_values3 = [headers3] + rows3
        self.table3 = CTkTable(
            master=self.scroll_frame, 
            row=len(table_values3), 
            column=len(table_values3[0]), 
            values=table_values3,
            header_color="#1f538d",
            colors=["#2a2a2a", "#242424"],
            hover_color="#3a3a3a",
        )
        self.table3.pack(expand=True, fill="both", padx=10, pady=10)

        # Dibujar Tabla 4
        if df_sim is not None:
            self.lbl_table4 = ctk.CTkLabel(self.scroll_frame, text="Simulación de Producción Anual", font=("Roboto", 16, "bold"), text_color="#3484F0")
            self.lbl_table4.pack(pady=(30, 5))
            headers4 = list(df_sim.columns)
            rows4 = df_sim.to_numpy().tolist()
            table_values4 = [headers4] + rows4
            self.table4 = CTkTable(
                master=self.scroll_frame, 
                row=len(table_values4), 
                column=len(table_values4[0]), 
                values=table_values4,
                header_color="#1f538d",
                colors=["#2a2a2a", "#242424"],
                hover_color="#3a3a3a",
            )
            self.table4.pack(expand=True, fill="both", padx=10, pady=20)

    def on_load_error(self, err_msg, error_details):
        self.set_loading_state(False)
        self.lbl_file.configure(text="Error al cargar")
        
        try:
            print(colored(f"\n[!] ERROR EN HILO DE CARGA:\n{err_msg}\n\nTRACEBACK:\n{error_details}", "red"))
        except:
            print(f"\n[!] ERROR EN HILO DE CARGA:\n{err_msg}\n\nTRACEBACK:\n{error_details}")
            
        messagebox.showerror("Error Detallado", f"Ocurrió un error al leer el archivo:\n{err_msg}\n\nDetalles técnicos:\n{error_details}")

    def send_to_powerpoint(self):
        if self.df_data is None or self.df_snr is None or self.df_prod is None:
            messagebox.showerror("Error", "Primero debes buscar un archivo Excel para extraer los datos.")
            return

        file_path = filedialog.askopenfilename(
            title="Seleccionar plantilla PowerPoint",
            initialdir=self.default_dir,
            filetypes=[("Archivos PowerPoint", "*.pptx")]
        )

        if file_path:
            save_path = filedialog.asksaveasfilename(
                title="Guardar presentación actualizada",
                initialdir=self.default_dir,
                initialfile="Proceso y Producciones_Actualizado.pptx",
                defaultextension=".pptx",
                filetypes=[("Archivos PowerPoint", "*.pptx")]
            )

            if save_path:
                self.set_loading_state(True, "Inyectando datos a PowerPoint...")
                threading.Thread(target=self.async_send_to_pptx, args=(file_path, save_path), daemon=True).start()

    def update_slide_chart(self, chart, categories, proceso_vals, diario_vals, programa_vals, columna1_vals, wine_color, green_color):
        from pptx.chart.data import CategoryChartData
        from pptx.util import Pt

        # Enviar a la gráfica valores ya redondeados a enteros.
        # (los cálculos internos más abajo siguen usando los datos originales)
        def _round(v):
            if v is None:
                return None
            try:
                return int(round(float(v)))
            except:
                return v

        proceso_r = [_round(v) for v in proceso_vals]
        diario_r = [_round(v) for v in diario_vals]
        programa_r = [_round(v) for v in programa_vals]
        columna1_r = [_round(v) for v in columna1_vals]

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('PROCESO', tuple(proceso_r))
        chart_data.add_series('Diario', tuple(diario_r))
        chart_data.add_series('Programa', tuple(programa_r))
        chart_data.add_series('Columna1', tuple(columna1_r))

        # Reemplazar los datos de la gráfica
        chart.replace_data(chart_data)

        # --- Restaurar y pintar los colores dinámicamente ---
        if len(chart.series) > 0:
            series = chart.series[0]
            
            # Buscar el último mes con producción real en la nueva lista de categorías
            last_month_idx = -1
            for i in range(17):
                if i < len(categories):
                    cat = categories[i]
                    val = proceso_vals[i]
                    # Si es un mes (contiene letras)
                    if any(c.isalpha() for c in cat):
                        if val is not None and val != 0:
                            last_month_idx = i

            # Eliminar TODOS los formatos específicos de puntos (<c:dPt>)
            # Esto obliga a que todas las barras hereden el color gris por defecto de la serie
            try:
                ser_el = series._element
                ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
                
                for _ in range(2):
                    dpts = ser_el.findall('c:dPt', ns)
                    for dpt in dpts:
                        ser_el.remove(dpt)
            except Exception:
                pass

            # Definir color gris para los años
            from pptx.dml.color import RGBColor
            gray_color = RGBColor(192, 192, 192)

            # Pintar todas las barras
            for p_idx in range(min(17, len(series.points))):
                try:
                    cat = categories[p_idx]
                    point = series.points[p_idx]
                    fill = point.format.fill
                    fill.solid()
                    
                    # Si es un mes, aplicamos vino o verde
                    if any(c.isalpha() for c in cat):
                        if p_idx == last_month_idx:
                            fill.fore_color.rgb = wine_color
                            try:
                                point.data_label.font.size = Pt(14)
                                point.data_label.font.bold = True
                            except Exception:
                                pass
                        else:
                            fill.fore_color.rgb = green_color
                            
                        # Limpiar modificadores de color
                        try:
                            ns_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            srgb_clr = point._element.find('.//a:srgbClr', ns_a)
                            if srgb_clr is not None:
                                for child in list(srgb_clr):
                                    if 'lumMod' in child.tag or 'lumOff' in child.tag:
                                        srgb_clr.remove(child)
                        except Exception:
                            pass
                    else:
                        # Si es un año, aplicamos el gris explícitamente
                        fill.fore_color.rgb = gray_color
                        
                except Exception:
                    pass

    def async_send_to_pptx(self, file_path, save_path):
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor

            prs = Presentation(file_path)
            if len(prs.slides) < 12:
                raise ValueError("La presentación debe tener al menos 12 diapositivas (incluyendo las de Cadereyta -Crudo, Gasolinas y Diesel).")
 
            # --- 1. PROCESAR DIAPOSITIVA DE CRUDO (DIAPOSITIVA 2) ---
            slide = prs.slides[1]
            chart = None

            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    break

            if not chart:
                raise ValueError("No se encontró ninguna gráfica en la segunda diapositiva (Crudo).")

            # Columnas específicas de la tabla 1 para Crudo: 'Crudo' (0), 'Cadereyta' (1)
            # En la tabla, "Cadereyta" es la que va en la columna de proceso.
            # Los datos de producción diaria están en la segunda columna (df_data.iloc[:, 1])
            proceso_col = self.df_data.columns[1]

            categories = []
            proceso_vals = []
            diario_vals = []
            programa_vals = []
            columna1_vals = []

            # Filtrar df_prod para quedarnos con los años y los meses activos
            prod_rows = []
            for idx, row in self.df_prod.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]
                if not cat: continue
                if not any(c.isalpha() for c in cat):
                    prod_rows.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                        if p_val != 0: prod_rows.append((cat, val))
                    except: pass
            
            if len(prod_rows) > 30: prod_rows = prod_rows[-30:]

            for i in range(len(prod_rows)):
                categories.append(prod_rows[i][0])
                try: proceso_vals.append(float(prod_rows[i][1]))
                except: proceso_vals.append(None)
                diario_vals.append(None)
                programa_vals.append(None)
                columna1_vals.append(None)

            # Llenar datos diarios (31 días)
            for i in range(31):
                categories.append(str(i + 1))
                proceso_vals.append(None)
                
                try: diario_vals.append(float(self.df_data[proceso_col].iloc[i]))
                except: diario_vals.append(None)
                
                try: programa_vals.append(float(self.df_snr.iloc[i, 0]))
                except: programa_vals.append(None)
                
                try: columna1_vals.append(float(self.df_snr.iloc[i, 1]))
                except: columna1_vals.append(None)

            # --- EXTRAER COLORES DE LA PLANTILLA ---
            wine_color = RGBColor(0x69, 0x19, 0x32)
            green_color = RGBColor(0x24, 0x5C, 0x4F)

            # Actualizar gráfica
            self.update_slide_chart(chart, categories, proceso_vals, diario_vals, programa_vals, columna1_vals, wine_color, green_color)


            # --- NUEVA SECCIÓN: PROCESAR DIAPOSITIVA 10 (Crudo Cadereyta) ---
            slide_cad = prs.slides[9]
            chart_cad = None
            for shape in slide_cad.shapes:
                if shape.has_chart:
                    chart_cad = shape.chart
                    break
            
            if chart_cad:
                categories_cad = []
                proceso_vals_cad = []
                diario_vals_cad = []
                programa_vals_cad = []
                columna1_vals_cad = []

                # Filtrar df_prod_cad
                prod_rows_cad = []
                for idx, row in self.df_prod_cad.iterrows():
                    cat = str(row.iloc[0]).strip()
                    val = row.iloc[1]
                    if not cat: continue
                    if not any(c.isalpha() for c in cat):
                        prod_rows_cad.append((cat, val))
                    else:
                        try:
                            p_val = float(val)
                            if p_val != 0: prod_rows_cad.append((cat, val))
                        except: pass
                
                if len(prod_rows_cad) > 30: prod_rows_cad = prod_rows_cad[-30:]
                
                for i in range(len(prod_rows_cad)):
                    categories_cad.append(prod_rows_cad[i][0])
                    try: proceso_vals_cad.append(float(prod_rows_cad[i][1]))
                    except: proceso_vals_cad.append(None)
                    diario_vals_cad.append(None)
                    programa_vals_cad.append(None)
                    columna1_vals_cad.append(None)

                proceso_col_cad = self.df_data_cad.columns[1]

                for i in range(31):
                    categories_cad.append(str(i + 1))
                    proceso_vals_cad.append(None)
                    
                    try: diario_vals_cad.append(float(self.df_data_cad[proceso_col_cad].iloc[i]))
                    except: diario_vals_cad.append(None)
                    
                    try: programa_vals_cad.append(float(self.df_snr_cad.iloc[i, 0]))
                    except: programa_vals_cad.append(None)
                    
                    try: columna1_vals_cad.append(float(self.df_snr_cad.iloc[i, 1]))
                    except: columna1_vals_cad.append(None)

                self.update_slide_chart(chart_cad, categories_cad, proceso_vals_cad, diario_vals_cad, programa_vals_cad, columna1_vals_cad, wine_color, green_color)


            # --- 2. PROCESAR DIAPOSITIVA DE GASOLINAS (DIAPOSITIVA 3) ---
            slide_gas = prs.slides[2]
            chart_gas = None
            for shape in slide_gas.shapes:
                if shape.has_chart:
                    chart_gas = shape.chart
                    break

            if not chart_gas:
                raise ValueError("No se encontró ninguna gráfica en la tercera diapositiva (Gasolinas).")

            snr_col_gas = None
            for col in self.df_data_gasolinas.columns:
                if "SNR" in str(col).upper():
                    snr_col_gas = col
                    break

            if not snr_col_gas:
                raise ValueError("No se encontró la columna 'SNR' en la tabla de Gasolinas.")

            categories_gas = []
            proceso_vals_gas = []
            diario_vals_gas = []
            programa_vals_gas = []
            columna1_vals_gas = []

            # Filtrar df_prod_gasolinas para quedarnos con los años y meses activos
            prod_rows_gas = []
            for idx, row in self.df_prod_gasolinas.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]
                
                if not cat:
                    continue
                    
                if not any(c.isalpha() for c in cat):
                    prod_rows_gas.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                    except:
                        p_val = 0
                    if p_val != 0:
                        prod_rows_gas.append((cat, val))

            # Ajustar al límite de 30 categorías
            if len(prod_rows_gas) > 30:
                prod_rows_gas = prod_rows_gas[-30:]

            # Llenar filas de años y meses (sin celdas vacías al final)
            for i in range(len(prod_rows_gas)):
                cat_val = prod_rows_gas[i][0]
                try:
                    proc_val = float(prod_rows_gas[i][1])
                except:
                    proc_val = None
                        
                categories_gas.append(cat_val)
                proceso_vals_gas.append(proc_val)
                diario_vals_gas.append(None)
                programa_vals_gas.append(None)
                columna1_vals_gas.append(None)

            # Llenar filas diarias
            for i in range(31):
                categories_gas.append(str(i + 1))
                proceso_vals_gas.append(None)
                
                d_val = None
                if i < len(self.df_data_gasolinas):
                    try:
                        d_val = float(self.df_data_gasolinas[snr_col_gas].iloc[i])
                    except:
                        d_val = None
                diario_vals_gas.append(d_val)
                
                p_val = None
                if i < len(self.df_snr_gasolinas):
                    try:
                        p_val = float(self.df_snr_gasolinas.iloc[i, 0])
                    except:
                        p_val = None
                programa_vals_gas.append(p_val)
                
                c_val = None
                if i < len(self.df_snr_gasolinas):
                    try:
                        c_val = float(self.df_snr_gasolinas.iloc[i, 1])
                    except:
                        c_val = None
                columna1_vals_gas.append(c_val)

            # Actualizar la gráfica de Gasolinas (Diapositiva 3)
            self.update_slide_chart(chart_gas, categories_gas, proceso_vals_gas, diario_vals_gas, programa_vals_gas, columna1_vals_gas, wine_color, green_color)


            # --- 3. PROCESAR DIAPOSITIVA DE DIESEL (DIAPOSITIVA 4) ---
            slide_die = prs.slides[3]
            chart_die = None
            for shape in slide_die.shapes:
                if shape.has_chart:
                    chart_die = shape.chart
                    break

            if not chart_die:
                raise ValueError("No se encontró ninguna gráfica en la cuarta diapositiva (Diesel).")

            snr_col_die = None
            for col in self.df_data_diesel.columns:
                if "SNR" in str(col).upper():
                    snr_col_die = col
                    break

            if not snr_col_die:
                raise ValueError("No se encontró la columna 'SNR' en la tabla de Diesel.")

            categories_die = []
            proceso_vals_die = []
            diario_vals_die = []
            programa_vals_die = []
            columna1_vals_die = []

            # Filtrar df_prod_diesel para quedarnos con los años y meses activos
            prod_rows_die = []
            for idx, row in self.df_prod_diesel.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]
                
                if not cat:
                    continue
                    
                if not any(c.isalpha() for c in cat):
                    prod_rows_die.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                    except:
                        p_val = 0
                    if p_val != 0:
                        prod_rows_die.append((cat, val))

            # Ajustar al límite de 30 categorías
            if len(prod_rows_die) > 30:
                prod_rows_die = prod_rows_die[-30:]

            # Llenar filas de años y meses (sin celdas vacías al final)
            for i in range(len(prod_rows_die)):
                cat_val = prod_rows_die[i][0]
                try:
                    proc_val = float(prod_rows_die[i][1])
                except:
                    proc_val = None
                        
                categories_die.append(cat_val)
                proceso_vals_die.append(proc_val)
                diario_vals_die.append(None)
                programa_vals_die.append(None)
                columna1_vals_die.append(None)

            # Llenar filas diarias
            for i in range(31):
                categories_die.append(str(i + 1))
                proceso_vals_die.append(None)
                
                d_val = None
                if i < len(self.df_data_diesel):
                    try:
                        d_val = float(self.df_data_diesel[snr_col_die].iloc[i])
                    except:
                        d_val = None
                diario_vals_die.append(d_val)
                
                p_val = None
                if i < len(self.df_snr_diesel):
                    try:
                        p_val = float(self.df_snr_diesel.iloc[i, 0])
                    except:
                        p_val = None
                programa_vals_die.append(p_val)
                
                c_val = None
                if i < len(self.df_snr_diesel):
                    try:
                        c_val = float(self.df_snr_diesel.iloc[i, 1])
                    except:
                        c_val = None
                columna1_vals_die.append(c_val)

            # Actualizar la gráfica de Diesel (Diapositiva 4)
            self.update_slide_chart(chart_die, categories_die, proceso_vals_die, diario_vals_die, programa_vals_die, columna1_vals_die, wine_color, green_color)


            # --- 4. PROCESAR DIAPOSITIVA DE TURBOSINA (DIAPOSITIVA 5) ---
            slide_turb = prs.slides[4]
            chart_turb = None
            for shape in slide_turb.shapes:
                if shape.has_chart:
                    chart_turb = shape.chart
                    break

            if not chart_turb:
                raise ValueError("No se encontró ninguna gráfica en la quinta diapositiva (Turbosina).")

            snr_col_turb = None
            for col in self.df_data_turbosina.columns:
                if "SNR" in str(col).upper():
                    snr_col_turb = col
                    break

            if not snr_col_turb:
                raise ValueError("No se encontró la columna 'SNR' en la tabla de Turbosina.")

            categories_turb = []
            proceso_vals_turb = []
            diario_vals_turb = []
            programa_vals_turb = []
            columna1_vals_turb = []

            # Filtrar df_prod_turbosina para quedarnos con los años y meses activos
            prod_rows_turb = []
            for idx, row in self.df_prod_turbosina.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]
                
                if not cat:
                    continue
                    
                if not any(c.isalpha() for c in cat):
                    prod_rows_turb.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                    except:
                        p_val = 0
                    if p_val != 0:
                        prod_rows_turb.append((cat, val))

            # Ajustar al límite de 30 categorías
            if len(prod_rows_turb) > 30:
                prod_rows_turb = prod_rows_turb[-30:]

            # Llenar filas de años y meses (sin celdas vacías al final)
            for i in range(len(prod_rows_turb)):
                cat_val = prod_rows_turb[i][0]
                try:
                    proc_val = float(prod_rows_turb[i][1])
                except:
                    proc_val = None
                        
                categories_turb.append(cat_val)
                proceso_vals_turb.append(proc_val)
                diario_vals_turb.append(None)
                programa_vals_turb.append(None)
                columna1_vals_turb.append(None)

            # Llenar filas diarias
            for i in range(31):
                categories_turb.append(str(i + 1))
                proceso_vals_turb.append(None)
                
                d_val = None
                if i < len(self.df_data_turbosina):
                    try:
                        d_val = float(self.df_data_turbosina[snr_col_turb].iloc[i])
                    except:
                        d_val = None
                diario_vals_turb.append(d_val)
                
                p_val = None
                if i < len(self.df_snr_turbosina):
                    try:
                        p_val = float(self.df_snr_turbosina.iloc[i, 0])
                    except:
                        p_val = None
                programa_vals_turb.append(p_val)
                
                c_val = None
                if i < len(self.df_snr_turbosina):
                    try:
                        c_val = float(self.df_snr_turbosina.iloc[i, 1])
                    except:
                        c_val = None
                columna1_vals_turb.append(c_val)

            # Actualizar la gráfica de Turbosina (Diapositiva 5)
            self.update_slide_chart(chart_turb, categories_turb, proceso_vals_turb, diario_vals_turb, programa_vals_turb, columna1_vals_turb, wine_color, green_color)


            # --- 5. PROCESAR DIAPOSITIVA DE ASFALTO (DIAPOSITIVA 6) ---
            if self.df_data_asfalto is not None and self.df_snr_asfalto is not None and self.df_prod_asfalto is not None:
                slide_asf = prs.slides[5]
                chart_asf = None
                for shape in slide_asf.shapes:
                    if shape.has_chart:
                        chart_asf = shape.chart
                        break

                if not chart_asf:
                    raise ValueError("No se encontró ninguna gráfica en la sexta diapositiva (Asfalto).")

                snr_col_asf = None
                for col in self.df_data_asfalto.columns:
                    if "SNR" in str(col).upper() or "REAL" in str(col).upper():
                        snr_col_asf = col
                        break
                if not snr_col_asf and len(self.df_data_asfalto.columns) > 1:
                    snr_col_asf = self.df_data_asfalto.columns[1]

                if not snr_col_asf:
                    raise ValueError("No se encontró la columna de producción diaria ('REAL') en la tabla de Asfalto.")

                categories_asf = []
                proceso_vals_asf = []
                diario_vals_asf = []
                programa_vals_asf = []
                columna1_vals_asf = []

                # Filtrar df_prod_asfalto para quedarnos con los años y meses activos
                prod_rows_asf = []
                for idx, row in self.df_prod_asfalto.iterrows():
                    cat = str(row.iloc[0]).strip()
                    val = row.iloc[1]

                    if not cat:
                        continue

                    if not any(c.isalpha() for c in cat):
                        prod_rows_asf.append((cat, val))
                    else:
                        try:
                            p_val = float(val)
                        except:
                            p_val = 0
                        if p_val != 0:
                            prod_rows_asf.append((cat, val))

                # Ajustar al límite de 30 categorías
                if len(prod_rows_asf) > 30:
                    prod_rows_asf = prod_rows_asf[-30:]

                # Llenar filas de años y meses (sin celdas vacías al final)
                for i in range(len(prod_rows_asf)):
                    cat_val = prod_rows_asf[i][0]
                    try:
                        proc_val = float(prod_rows_asf[i][1])
                    except:
                        proc_val = None

                    categories_asf.append(cat_val)
                    proceso_vals_asf.append(proc_val)
                    diario_vals_asf.append(None)
                    programa_vals_asf.append(None)
                    columna1_vals_asf.append(None)

                # Llenar filas diarias
                for i in range(31):
                    categories_asf.append(str(i + 1))
                    proceso_vals_asf.append(None)

                    d_val = None
                    if i < len(self.df_data_asfalto):
                        try:
                            d_val = float(self.df_data_asfalto[snr_col_asf].iloc[i])
                        except:
                            d_val = None
                    diario_vals_asf.append(d_val)

                    p_val = None
                    if i < len(self.df_snr_asfalto):
                        try:
                            p_val = float(self.df_snr_asfalto.iloc[i, 0])
                        except:
                            p_val = None
                    programa_vals_asf.append(p_val)

                    c_val = None
                    if i < len(self.df_snr_asfalto):
                        try:
                            c_val = float(self.df_snr_asfalto.iloc[i, 1])
                        except:
                            c_val = None
                    columna1_vals_asf.append(c_val)

                # Actualizar la gráfica de Asfalto (Diapositiva 6)
                self.update_slide_chart(chart_asf, categories_asf, proceso_vals_asf, diario_vals_asf, programa_vals_asf, columna1_vals_asf, wine_color, green_color)


            # --- 6. PROCESAR DIAPOSITIVA DE COMBUSTOLEO (DIAPOSITIVA 7) ---
            if self.df_data_combustoleo is not None and self.df_snr_combustoleo is not None and self.df_prod_combustoleo is not None:
                slide_comb = prs.slides[6]
                chart_comb = None
                for shape in slide_comb.shapes:
                    if shape.has_chart:
                        chart_comb = shape.chart
                        break

                if not chart_comb:
                    raise ValueError("No se encontró ninguna gráfica en la séptima diapositiva (Combustoleo).")

                # La columna de producción diaria real es "Combustoleo" (no la de SNR)
                diario_col_comb = None
                for col in self.df_data_combustoleo.columns:
                    if "SNR" not in str(col).upper() and "REAL" not in str(col).upper():
                        diario_col_comb = col
                        break
                if not diario_col_comb and len(self.df_data_combustoleo.columns) > 1:
                    diario_col_comb = self.df_data_combustoleo.columns[1]
                elif not diario_col_comb:
                    diario_col_comb = self.df_data_combustoleo.columns[0]

                if not diario_col_comb:
                    raise ValueError("No se encontró la columna de producción diaria ('Combustoleo') en la tabla de Combustoleo.")

                categories_comb = []
                proceso_vals_comb = []
                diario_vals_comb = []
                programa_vals_comb = []
                columna1_vals_comb = []

                # Filtrar df_prod_combustoleo para quedarnos con los años y meses activos
                prod_rows_comb = []
                for idx, row in self.df_prod_combustoleo.iterrows():
                    cat = str(row.iloc[0]).strip()
                    val = row.iloc[1]

                    if not cat:
                        continue

                    if not any(c.isalpha() for c in cat):
                        prod_rows_comb.append((cat, val))
                    else:
                        try:
                            p_val = float(val)
                        except:
                            p_val = 0
                        if p_val != 0:
                            prod_rows_comb.append((cat, val))

                # Ajustar al límite de 30 categorías
                if len(prod_rows_comb) > 30:
                    prod_rows_comb = prod_rows_comb[-30:]

                # Llenar filas de años y meses (sin celdas vacías al final)
                for i in range(len(prod_rows_comb)):
                    cat_val = prod_rows_comb[i][0]
                    try:
                        proc_val = float(prod_rows_comb[i][1])
                    except:
                        proc_val = None

                    categories_comb.append(cat_val)
                    proceso_vals_comb.append(proc_val)
                    diario_vals_comb.append(None)
                    programa_vals_comb.append(None)
                    columna1_vals_comb.append(None)

                # Llenar filas diarias
                for i in range(31):
                    categories_comb.append(str(i + 1))
                    proceso_vals_comb.append(None)

                    d_val = None
                    if i < len(self.df_data_combustoleo):
                        try:
                            d_val = float(self.df_data_combustoleo[diario_col_comb].iloc[i])
                        except:
                            d_val = None
                    diario_vals_comb.append(d_val)

                    p_val = None
                    if i < len(self.df_snr_combustoleo):
                        try:
                            p_val = float(self.df_snr_combustoleo.iloc[i, 0])
                        except:
                            p_val = None
                    programa_vals_comb.append(p_val)

                    c_val = None
                    if i < len(self.df_snr_combustoleo):
                        try:
                            c_val = float(self.df_snr_combustoleo.iloc[i, 1])
                        except:
                            c_val = None
                    columna1_vals_comb.append(c_val)

                # Actualizar la gráfica de Combustoleo (Diapositiva 7)
                self.update_slide_chart(chart_comb, categories_comb, proceso_vals_comb, diario_vals_comb, programa_vals_comb, columna1_vals_comb, wine_color, green_color)
 
                # --- 7. PROCESAR DIAPOSITIVA DE GASOLINAS CADEREYTA (DIAPOSITIVA 11) ---
                if self.df_data_cad_gas is not None and self.df_snr_cad_gas is not None and self.df_prod_cad_gas is not None:
                    slide_cad_gas = prs.slides[10]
                    chart_cad_gas = None
                    for shape in slide_cad_gas.shapes:
                        if shape.has_chart:
                            chart_cad_gas = shape.chart
                            break
                    
                    if chart_cad_gas:
                        # Buscar columna SNR o usar la segunda por defecto
                        snr_col_cad_gas = None
                        for col in self.df_data_cad_gas.columns:
                            if "SNR" in str(col).upper():
                                snr_col_cad_gas = col
                                break
                        if not snr_col_cad_gas and len(self.df_data_cad_gas.columns) >= 2:
                            snr_col_cad_gas = self.df_data_cad_gas.columns[1]
                        
                        if snr_col_cad_gas:
                            categories_cg = []
                            proceso_vals_cg = []
                            diario_vals_cg = []
                            programa_vals_cg = []
                            columna1_vals_cg = []

                            # Filtrar producción anual (últimas 30 categorías)
                            prod_rows_cg = []
                            for idx, row in self.df_prod_cad_gas.iterrows():
                                cat = str(row.iloc[0]).strip()
                                val = row.iloc[1]
                                if not cat: continue
                                if not any(c.isalpha() for c in cat):
                                    prod_rows_cg.append((cat, val))
                                else:
                                    try:
                                        if float(val) != 0: prod_rows_cg.append((cat, val))
                                    except: pass





                            if len(prod_rows_cg) > 30: prod_rows_cg = prod_rows_cg[-30:]

                            for i in range(len(prod_rows_cg)):
                                categories_cg.append(prod_rows_cg[i][0])
                                try: proceso_vals_cg.append(float(prod_rows_cg[i][1]))
                                except: proceso_vals_cg.append(None)
                                diario_vals_cg.append(None)
                                programa_vals_cg.append(None)
                                columna1_vals_cg.append(None)
 
                            # Llenar datos diarios (31 días)
                            for i in range(31):
                                categories_cg.append(str(i + 1))
                                proceso_vals_cg.append(None)
                                
                                try: diario_vals_cg.append(float(self.df_data_cad_gas[snr_col_cad_gas].iloc[i]))
                                except: diario_vals_cg.append(None)
                                
                                try: programa_vals_cg.append(float(self.df_snr_cad_gas.iloc[i, 0]))
                                except: programa_vals_cg.append(None)
                                
                                try: columna1_vals_cg.append(float(self.df_snr_cad_gas.iloc[i, 1]))
                                except: columna1_vals_cg.append(None)
 
                            self.update_slide_chart(chart_cad_gas, categories_cg, proceso_vals_cg, diario_vals_cg, programa_vals_cg, columna1_vals_cg, wine_color, green_color)
 
                    # --- 8. PROCESAR DIAPOSITIVA DE DIESEL CADEREYTA (DIAPOSITIVA 12) ---
                    if self.df_data_cad_die is not None and self.df_snr_cad_die is not None and self.df_prod_cad_die is not None:
                        slide_cad_die = prs.slides[11]
                        chart_cad_die = None
                        for shape in slide_cad_die.shapes:
                            if shape.has_chart:
                                chart_cad_die = shape.chart
                                break
                        
                        if chart_cad_die:
                            snr_col_cad_die = None
                            for col in self.df_data_cad_die.columns:
                                if "SNR" in str(col).upper():
                                    snr_col_cad_die = col
                                    break
                            if not snr_col_cad_die and len(self.df_data_cad_die.columns) >= 2:
                                snr_col_cad_die = self.df_data_cad_die.columns[1]
                            
                            if snr_col_cad_die:
                                categories_cd = []
                                proceso_vals_cd = []
                                diario_vals_cd = []
                                programa_vals_cd = []
                                columna1_vals_cd = []
 
                                prod_rows_cd = []
                                for idx, row in self.df_prod_cad_die.iterrows():
                                    cat = str(row.iloc[0]).strip()
                                    val = row.iloc[1]
                                    if not cat: continue
                                    if not any(c.isalpha() for c in cat):
                                        prod_rows_cd.append((cat, val))
                                    else:
                                        try:
                                            if float(val) != 0: prod_rows_cd.append((cat, val))
                                        except: pass
                                
                                if len(prod_rows_cd) > 30: prod_rows_cd = prod_rows_cd[-30:]
 
                                for i in range(len(prod_rows_cd)):
                                    categories_cd.append(prod_rows_cd[i][0])
                                    try: proceso_vals_cd.append(float(prod_rows_cd[i][1]))
                                    except: proceso_vals_cd.append(None)
                                    diario_vals_cd.append(None)
                                    programa_vals_cd.append(None)
                                    columna1_vals_cd.append(None)
 
                                # Llenar datos diarios (31 días)
                                for i in range(31):
                                    categories_cd.append(str(i + 1))
                                    proceso_vals_cd.append(None)
                                    
                                    try: diario_vals_cd.append(float(self.df_data_cad_die[snr_col_cad_die].iloc[i]))
                                    except: diario_vals_cd.append(None)
                                    
                                    try: programa_vals_cd.append(float(self.df_snr_cad_die.iloc[i, 0]))
                                    except: programa_vals_cd.append(None)
                                    
                                    try: columna1_vals_cd.append(float(self.df_snr_cad_die.iloc[i, 1]))
                                    except: columna1_vals_cd.append(None)
 
                                self.update_slide_chart(chart_cad_die, categories_cd, proceso_vals_cd, diario_vals_cd, programa_vals_cd, columna1_vals_cd, wine_color, green_color)

 
            prs.save(save_path)


            self.after(0, self.on_pptx_success, save_path)

        except Exception as e:
            err_details = traceback.format_exc()
            self.after(0, self.on_pptx_error, str(e), err_details)

    def on_pptx_success(self, save_path):
        self.set_loading_state(False)
        self.lbl_file.configure(text="Inyección a PowerPoint lista")
        messagebox.showinfo("Éxito", f"¡Datos de la gráfica actualizados con éxito!\n\nGuardado en:\n{save_path}")

    def on_pptx_error(self, err_msg, error_details):
        self.set_loading_state(False)
        self.lbl_file.configure(text="Error PowerPoint")
        
        try:
            print(colored(f"\n[!] ERROR EN HILO DE POWERPOINT:\n{err_msg}\n\nTRACEBACK:\n{error_details}", "red"))
        except:
            print(f"\n[!] ERROR EN HILO DE POWERPOINT:\n{err_msg}\n\nTRACEBACK:\n{error_details}")
            
        messagebox.showerror("Error", f"Ocurrió un error al procesar la presentación:\n{err_msg}")

if __name__ == "__main__":
    import signal
    import sys

    def sigint_handler(sig, frame):
        print("\nEjecución cancelada por el usuario (Ctrl+C). Saliendo...")
        sys.exit(0)

    signal.signal(signal.SIGINT, sigint_handler)

    app = ExcelViewerApp()

    def check_signals():
        app.after(500, check_signals)
        
    app.after(500, check_signals)
    app.mainloop()
