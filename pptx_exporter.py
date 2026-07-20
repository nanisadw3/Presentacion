import os
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

import traceback
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor

def update_slide_chart(chart, categories, proceso_vals, diario_vals, programa_vals, columna1_vals, wine_color, green_color, current_month_name=None):
    from pptx.chart.data import CategoryChartData
    from pptx.util import Pt

    # Enviar a la gráfica valores ya redondeados a enteros.
    # (los cálculos internos más abajo siguen usando los datos originales)
    def _round(v):
        if v is None:
            return None
        try:
            return int(round(float(v)))
        except:
            return v

    proceso_r = [_round(v) for v in proceso_vals]
    diario_r = [_round(v) for v in diario_vals]
    programa_r = [_round(v) for v in programa_vals]
    columna1_r = [_round(v) for v in columna1_vals]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('PROCESO', tuple(proceso_r))
    chart_data.add_series('Diario', tuple(diario_r))
    chart_data.add_series('Programa', tuple(programa_r))
    chart_data.add_series('Columna1', tuple(columna1_r))

    # Reemplazar los datos de la gráfica
    chart.replace_data(chart_data)

    # --- Restaurar y pintar los colores dinámicamente ---
    if len(chart.series) > 0:
        series = chart.series[0]
        
        # Buscar el último mes con producción real en la nueva lista de categorías
        # Determinar cuál es el último mes con datos para pintarlo de vino
        last_month_idx = -1
        for i, val in enumerate(proceso_vals):
            if i < len(categories) and any(c.isalpha() for c in categories[i]):
                if current_month_name and categories[i].strip().lower()[:3] == current_month_name.strip().lower()[:3]:
                    last_month_idx = i
                elif val is not None and val != 0:
                    last_month_idx = i

        # Eliminar TODOS los formatos específicos de puntos (<c:dPt>)
        # Esto obliga a que todas las barras hereden el color gris por defecto de la serie
        try:
            ser_el = series._element
            ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
            
            for _ in range(2):
                dpts = ser_el.findall('c:dPt', ns)
                for dpt in dpts:
                    ser_el.remove(dpt)
        except Exception:
            pass

        # Definir color gris para los años
        from pptx.dml.color import RGBColor
        gray_color = RGBColor(192, 192, 192)

        # Pintar todas las barras y formatear sus etiquetas de datos
        for p_idx in range(min(17, len(series.points))):
            try:
                cat = categories[p_idx]
                point = series.points[p_idx]
                fill = point.format.fill
                fill.solid()
                
                # Si es un mes, aplicamos vino o verde
                if any(c.isalpha() for c in cat):
                    if p_idx == last_month_idx:
                        fill.fore_color.rgb = wine_color
                        try:
                            point.data_label.font.size = Pt(14)
                            point.data_label.font.bold = True
                        except Exception:
                            pass
                    else:
                        fill.fore_color.rgb = green_color
                        try:
                            point.data_label.font.size = Pt(9)
                            point.data_label.font.bold = False
                        except Exception:
                            pass
                        
                    # Limpiar modificadores de color
                    try:
                        ns_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        srgb_clr = point._element.find('.//a:srgbClr', ns_a)
                        if srgb_clr is not None:
                            for child in list(srgb_clr):
                                if 'lumMod' in child.tag or 'lumOff' in child.tag:
                                    srgb_clr.remove(child)
                    except Exception:
                        pass
                else:
                    # Si es un año, aplicamos el gris explícitamente y reseteamos fuente
                    fill.fore_color.rgb = gray_color
                    try:
                        point.data_label.font.size = Pt(9)
                        point.data_label.font.bold = False
                    except Exception:
                        pass
                    
            except Exception:
                pass


def export_to_pptx(app, file_path, save_path):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        prs = Presentation(file_path)
        if len(prs.slides) < 52:
            raise ValueError("La presentación debe tener al menos 52 diapositivas (incluyendo las de Madero, Minatitlán, Salamanca, Salina Cruz, Tula y Olmeca).")
 
        # --- 1. PROCESAR DIAPOSITIVA DE CRUDO (DIAPOSITIVA 2) ---
        slide = prs.slides[1]
        chart = None

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
            elif shape.has_text_frame and "CMP:" in shape.text:
                cmp_val = getattr(app, 'cmp_value', "1234.8")
                import re
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "CMP:" in run.text:
                            run.text = re.sub(r'CMP:\s*[\d\.]+', f'CMP: {cmp_val}', run.text)

        if not chart:
            raise ValueError("No se encontró ninguna gráfica en la segunda diapositiva (Crudo).")

        # Columnas específicas de la tabla 1 para Crudo: 'Crudo' (0), 'Cadereyta' (1)
        # En la tabla, "Cadereyta" es la que va en la columna de proceso.
        # Los datos de producción diaria están en la segunda columna (df_data.iloc[:, 1])
        proceso_col = app.df_data.columns[1]

        categories = []
        proceso_vals = []
        diario_vals = []
        programa_vals = []
        columna1_vals = []

        # Filtrar df_prod para quedarnos con los años y los meses activos
        prod_rows = []
        for idx, row in app.df_prod.iterrows():
            cat = str(row.iloc[0]).strip()
            val = row.iloc[1]
            if not cat: continue
            if not any(c.isalpha() for c in cat):
                prod_rows.append((cat, val))
            else:
                try:
                    p_val = float(val)
                    if p_val != 0: prod_rows.append((cat, val))
                except: pass
        
        if len(prod_rows) > 30: prod_rows = prod_rows[-30:]

        master_current_month = None
        for cat, val in prod_rows:
            if any(c.isalpha() for c in cat):
                master_current_month = cat

        for i in range(len(prod_rows)):
            categories.append(prod_rows[i][0])
            try: proceso_vals.append(float(prod_rows[i][1]))
            except: proceso_vals.append(None)
            diario_vals.append(None)
            programa_vals.append(None)
            columna1_vals.append(None)

        # Llenar datos diarios (31 días)
        for i in range(31):
            categories.append(str(i + 1))
            proceso_vals.append(None)
            
            try: diario_vals.append(float(app.df_data[proceso_col].iloc[i]))
            except: diario_vals.append(None)
            
            try: programa_vals.append(float(app.df_snr.iloc[i, 0]))
            except: programa_vals.append(None)
            
            try: columna1_vals.append(float(app.df_snr.iloc[i, 1]))
            except: columna1_vals.append(None)

        # --- EXTRAER COLORES DE LA PLANTILLA ---
        wine_color = RGBColor(0x69, 0x19, 0x32)
        green_color = RGBColor(0x24, 0x5C, 0x4F)

        # Actualizar gráfica
        update_slide_chart(chart, categories, proceso_vals, diario_vals, programa_vals, columna1_vals, wine_color, green_color)


        # --- NUEVA SECCIÓN: PROCESAR DIAPOSITIVA 10 (Crudo Cadereyta) ---
        slide_cad = prs.slides[9]
        chart_cad = None
        for shape in slide_cad.shapes:
            if shape.has_chart:
                chart_cad = shape.chart
                break
        
        if chart_cad:
            categories_cad = []
            proceso_vals_cad = []
            diario_vals_cad = []
            programa_vals_cad = []
            columna1_vals_cad = []

            # Filtrar df_prod_cad
            prod_rows_cad = []
            for idx, row in app.df_prod_cad.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]
                if not cat: continue
                if not any(c.isalpha() for c in cat):
                    prod_rows_cad.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                        if p_val != 0: prod_rows_cad.append((cat, val))
                    except: pass
            
            if len(prod_rows_cad) > 30: prod_rows_cad = prod_rows_cad[-30:]
            
            for i in range(len(prod_rows_cad)):
                categories_cad.append(prod_rows_cad[i][0])
                try: proceso_vals_cad.append(float(prod_rows_cad[i][1]))
                except: proceso_vals_cad.append(None)
                diario_vals_cad.append(None)
                programa_vals_cad.append(None)
                columna1_vals_cad.append(None)

            proceso_col_cad = app.df_data_cad.columns[1]

            for i in range(31):
                categories_cad.append(str(i + 1))
                proceso_vals_cad.append(None)
                
                try: diario_vals_cad.append(float(app.df_data_cad[proceso_col_cad].iloc[i]))
                except: diario_vals_cad.append(None)
                
                try: programa_vals_cad.append(float(app.df_snr_cad.iloc[i, 0]))
                except: programa_vals_cad.append(None)
                
                try: columna1_vals_cad.append(float(app.df_snr_cad.iloc[i, 1]))
                except: columna1_vals_cad.append(None)

            update_slide_chart(chart_cad, categories_cad, proceso_vals_cad, diario_vals_cad, programa_vals_cad, columna1_vals_cad, wine_color, green_color)


        # --- 2. PROCESAR DIAPOSITIVA DE GASOLINAS (DIAPOSITIVA 3) ---
        slide_gas = prs.slides[2]
        chart_gas = None
        for shape in slide_gas.shapes:
            if shape.has_chart:
                chart_gas = shape.chart
            elif shape.has_text_frame and "CMP" in shape.text and "Demanda" in shape.text:
                cmp_val = getattr(app, 'cmp_gasolinas', "513")
                import re
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "CMP" in run.text:
                            run.text = re.sub(r'CMP\s*:\s*[\d\.]+', f'CMP : {cmp_val}', run.text)

        if not chart_gas:
            raise ValueError("No se encontró ninguna gráfica en la tercera diapositiva (Gasolinas).")

        snr_col_gas = None
        for col in app.df_data_gasolinas.columns:
            if "SNR" in str(col).upper():
                snr_col_gas = col
                break

        if not snr_col_gas:
            raise ValueError("No se encontró la columna 'SNR' en la tabla de Gasolinas.")

        categories_gas = []
        proceso_vals_gas = []
        diario_vals_gas = []
        programa_vals_gas = []
        columna1_vals_gas = []

        # Filtrar df_prod_gasolinas para quedarnos con los años y meses activos
        prod_rows_gas = []
        for idx, row in app.df_prod_gasolinas.iterrows():
            cat = str(row.iloc[0]).strip()
            val = row.iloc[1]
            
            if not cat:
                continue
                
            if not any(c.isalpha() for c in cat):
                prod_rows_gas.append((cat, val))
            else:
                try:
                    p_val = float(val)
                except:
                    p_val = 0
                if p_val != 0:
                    prod_rows_gas.append((cat, val))

        # Ajustar al límite de 30 categorías
        if len(prod_rows_gas) > 30:
            prod_rows_gas = prod_rows_gas[-30:]

        # Llenar filas de años y meses (sin celdas vacías al final)
        for i in range(len(prod_rows_gas)):
            cat_val = prod_rows_gas[i][0]
            try:
                proc_val = float(prod_rows_gas[i][1])
            except:
                proc_val = None
                    
            categories_gas.append(cat_val)
            proceso_vals_gas.append(proc_val)
            diario_vals_gas.append(None)
            programa_vals_gas.append(None)
            columna1_vals_gas.append(None)

        # Llenar filas diarias
        for i in range(31):
            categories_gas.append(str(i + 1))
            proceso_vals_gas.append(None)
            
            d_val = None
            if i < len(app.df_data_gasolinas):
                try:
                    d_val = float(app.df_data_gasolinas[snr_col_gas].iloc[i])
                except:
                    d_val = None
            diario_vals_gas.append(d_val)
            
            p_val = None
            if i < len(app.df_snr_gasolinas):
                try:
                    p_val = float(app.df_snr_gasolinas.iloc[i, 0])
                except:
                    p_val = None
            programa_vals_gas.append(p_val)
            
            c_val = None
            if i < len(app.df_snr_gasolinas):
                try:
                    c_val = float(app.df_snr_gasolinas.iloc[i, 1])
                except:
                    c_val = None
            columna1_vals_gas.append(c_val)

        # Actualizar la gráfica de Gasolinas (Diapositiva 3)
        update_slide_chart(chart_gas, categories_gas, proceso_vals_gas, diario_vals_gas, programa_vals_gas, columna1_vals_gas, wine_color, green_color)


        # --- 3. PROCESAR DIAPOSITIVA DE DIESEL (DIAPOSITIVA 4) ---
        slide_die = prs.slides[3]
        chart_die = None
        for shape in slide_die.shapes:
            if shape.has_chart:
                chart_die = shape.chart
            elif shape.has_text_frame and "CMP" in shape.text and "Demanda" in shape.text:
                cmp_val = getattr(app, 'cmp_diesel', "386.9")
                import re
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "CMP" in run.text:
                            run.text = re.sub(r'CMP\s*:\s*[\d\.]+', f'CMP : {cmp_val}', run.text)

        if not chart_die:
            raise ValueError("No se encontró ninguna gráfica en la cuarta diapositiva (Diesel).")

        snr_col_die = None
        for col in app.df_data_diesel.columns:
            if "SNR" in str(col).upper():
                snr_col_die = col
                break

        if not snr_col_die:
            raise ValueError("No se encontró la columna 'SNR' en la tabla de Diesel.")

        categories_die = []
        proceso_vals_die = []
        diario_vals_die = []
        programa_vals_die = []
        columna1_vals_die = []

        # Filtrar df_prod_diesel para quedarnos con los años y meses activos
        prod_rows_die = []
        for idx, row in app.df_prod_diesel.iterrows():
            cat = str(row.iloc[0]).strip()
            val = row.iloc[1]
            
            if not cat:
                continue
                
            if not any(c.isalpha() for c in cat):
                prod_rows_die.append((cat, val))
            else:
                try:
                    p_val = float(val)
                except:
                    p_val = 0
                if p_val != 0:
                    prod_rows_die.append((cat, val))

        # Ajustar al límite de 30 categorías
        if len(prod_rows_die) > 30:
            prod_rows_die = prod_rows_die[-30:]

        # Llenar filas de años y meses (sin celdas vacías al final)
        for i in range(len(prod_rows_die)):
            cat_val = prod_rows_die[i][0]
            try:
                proc_val = float(prod_rows_die[i][1])
            except:
                proc_val = None
                    
            categories_die.append(cat_val)
            proceso_vals_die.append(proc_val)
            diario_vals_die.append(None)
            programa_vals_die.append(None)
            columna1_vals_die.append(None)

        # Llenar filas diarias
        for i in range(31):
            categories_die.append(str(i + 1))
            proceso_vals_die.append(None)
            
            d_val = None
            if i < len(app.df_data_diesel):
                try:
                    d_val = float(app.df_data_diesel[snr_col_die].iloc[i])
                except:
                    d_val = None
            diario_vals_die.append(d_val)
            
            p_val = None
            if i < len(app.df_snr_diesel):
                try:
                    p_val = float(app.df_snr_diesel.iloc[i, 0])
                except:
                    p_val = None
            programa_vals_die.append(p_val)
            
            c_val = None
            if i < len(app.df_snr_diesel):
                try:
                    c_val = float(app.df_snr_diesel.iloc[i, 1])
                except:
                    c_val = None
            columna1_vals_die.append(c_val)

        # Actualizar la gráfica de Diesel (Diapositiva 4)
        update_slide_chart(chart_die, categories_die, proceso_vals_die, diario_vals_die, programa_vals_die, columna1_vals_die, wine_color, green_color)


        # --- 4. PROCESAR DIAPOSITIVA DE TURBOSINA (DIAPOSITIVA 5) ---
        slide_turb = prs.slides[4]
        chart_turb = None
        for shape in slide_turb.shapes:
            if shape.has_chart:
                chart_turb = shape.chart
            elif shape.has_text_frame and "CMP" in shape.text and "Demanda" in shape.text:
                cmp_val = getattr(app, 'cmp_turbosina', "312.4")
                import re
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "CMP" in run.text:
                            run.text = re.sub(r'CMP\s*:\s*[\d\.]+', f'CMP : {cmp_val}', run.text)

        if not chart_turb:
            raise ValueError("No se encontró ninguna gráfica en la quinta diapositiva (Turbosina).")

        snr_col_turb = None
        for col in app.df_data_turbosina.columns:
            if "SNR" in str(col).upper():
                snr_col_turb = col
                break

        if not snr_col_turb:
            raise ValueError("No se encontró la columna 'SNR' en la tabla de Turbosina.")

        categories_turb = []
        proceso_vals_turb = []
        diario_vals_turb = []
        programa_vals_turb = []
        columna1_vals_turb = []

        # Filtrar df_prod_turbosina para quedarnos con los años y meses activos
        prod_rows_turb = []
        for idx, row in app.df_prod_turbosina.iterrows():
            cat = str(row.iloc[0]).strip()
            val = row.iloc[1]
            
            if not cat:
                continue
                
            if not any(c.isalpha() for c in cat):
                prod_rows_turb.append((cat, val))
            else:
                try:
                    p_val = float(val)
                except:
                    p_val = 0
                if p_val != 0:
                    prod_rows_turb.append((cat, val))

        # Ajustar al límite de 30 categorías
        if len(prod_rows_turb) > 30:
            prod_rows_turb = prod_rows_turb[-30:]

        # Llenar filas de años y meses (sin celdas vacías al final)
        for i in range(len(prod_rows_turb)):
            cat_val = prod_rows_turb[i][0]
            try:
                proc_val = float(prod_rows_turb[i][1])
            except:
                proc_val = None
                    
            categories_turb.append(cat_val)
            proceso_vals_turb.append(proc_val)
            diario_vals_turb.append(None)
            programa_vals_turb.append(None)
            columna1_vals_turb.append(None)

        # Llenar filas diarias
        for i in range(31):
            categories_turb.append(str(i + 1))
            proceso_vals_turb.append(None)
            
            d_val = None
            if i < len(app.df_data_turbosina):
                try:
                    d_val = float(app.df_data_turbosina[snr_col_turb].iloc[i])
                except:
                    d_val = None
            diario_vals_turb.append(d_val)
            
            p_val = None
            if i < len(app.df_snr_turbosina):
                try:
                    p_val = float(app.df_snr_turbosina.iloc[i, 0])
                except:
                    p_val = None
            programa_vals_turb.append(p_val)
            
            c_val = None
            if i < len(app.df_snr_turbosina):
                try:
                    c_val = float(app.df_snr_turbosina.iloc[i, 1])
                except:
                    c_val = None
            columna1_vals_turb.append(c_val)

        # Actualizar la gráfica de Turbosina (Diapositiva 5)
        update_slide_chart(chart_turb, categories_turb, proceso_vals_turb, diario_vals_turb, programa_vals_turb, columna1_vals_turb, wine_color, green_color)


        # --- 5. PROCESAR DIAPOSITIVA DE ASFALTO (DIAPOSITIVA 6) ---
        if app.df_data_asfalto is not None and app.df_snr_asfalto is not None and app.df_prod_asfalto is not None:
            slide_asf = prs.slides[5]
            chart_asf = None
            for shape in slide_asf.shapes:
                if shape.has_chart:
                    chart_asf = shape.chart
                    break

            if not chart_asf:
                raise ValueError("No se encontró ninguna gráfica en la sexta diapositiva (Asfalto).")

            snr_col_asf = None
            for col in app.df_data_asfalto.columns:
                if "SNR" in str(col).upper() or "REAL" in str(col).upper():
                    snr_col_asf = col
                    break
            if not snr_col_asf and len(app.df_data_asfalto.columns) > 1:
                snr_col_asf = app.df_data_asfalto.columns[1]

            if not snr_col_asf:
                raise ValueError("No se encontró la columna de producción diaria ('REAL') en la tabla de Asfalto.")

            categories_asf = []
            proceso_vals_asf = []
            diario_vals_asf = []
            programa_vals_asf = []
            columna1_vals_asf = []

            # Filtrar df_prod_asfalto para quedarnos con los años y meses activos
            prod_rows_asf = []
            for idx, row in app.df_prod_asfalto.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]

                if not cat:
                    continue

                if not any(c.isalpha() for c in cat):
                    prod_rows_asf.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                    except:
                        p_val = 0
                    # Permitir meses con producción 0
                    prod_rows_asf.append((cat, val))

            # Ajustar al límite de 30 categorías
            if len(prod_rows_asf) > 30:
                prod_rows_asf = prod_rows_asf[-30:]

            # Llenar filas de años y meses (sin celdas vacías al final)
            for i in range(len(prod_rows_asf)):
                cat_val = prod_rows_asf[i][0]
                try:
                    proc_val = float(prod_rows_asf[i][1])
                except:
                    proc_val = None

                categories_asf.append(cat_val)
                proceso_vals_asf.append(proc_val)
                diario_vals_asf.append(None)
                programa_vals_asf.append(None)
                columna1_vals_asf.append(None)

            # Llenar filas diarias
            for i in range(31):
                categories_asf.append(str(i + 1))
                proceso_vals_asf.append(None)

                d_val = None
                if i < len(app.df_data_asfalto):
                    try:
                        d_val = float(app.df_data_asfalto[snr_col_asf].iloc[i])
                    except:
                        d_val = None
                diario_vals_asf.append(d_val)

                p_val = None
                if i < len(app.df_snr_asfalto):
                    try:
                        p_val = float(app.df_snr_asfalto.iloc[i, 0])
                    except:
                        p_val = None
                programa_vals_asf.append(p_val)

                c_val = None
                if i < len(app.df_snr_asfalto):
                    try:
                        c_val = float(app.df_snr_asfalto.iloc[i, 1])
                    except:
                        c_val = None
                columna1_vals_asf.append(c_val)

            # Actualizar la gráfica de Asfalto (Diapositiva 6)
            update_slide_chart(chart_asf, categories_asf, proceso_vals_asf, diario_vals_asf, programa_vals_asf, columna1_vals_asf, wine_color, green_color, current_month_name=master_current_month)


        # --- 6. PROCESAR DIAPOSITIVA DE COMBUSTOLEO (DIAPOSITIVA 7) ---
        if app.df_data_combustoleo is not None and app.df_snr_combustoleo is not None and app.df_prod_combustoleo is not None:
            slide_comb = prs.slides[6]
            chart_comb = None
            for shape in slide_comb.shapes:
                if shape.has_chart:
                    chart_comb = shape.chart
                    break

            if not chart_comb:
                raise ValueError("No se encontró ninguna gráfica en la séptima diapositiva (Combustoleo).")

            # La columna de producción diaria real es "Combustoleo" (no la de SNR)
            diario_col_comb = None
            for col in app.df_data_combustoleo.columns:
                if "SNR" not in str(col).upper() and "REAL" not in str(col).upper():
                    diario_col_comb = col
                    break
            if not diario_col_comb and len(app.df_data_combustoleo.columns) > 1:
                diario_col_comb = app.df_data_combustoleo.columns[1]
            elif not diario_col_comb:
                diario_col_comb = app.df_data_combustoleo.columns[0]

            if not diario_col_comb:
                raise ValueError("No se encontró la columna de producción diaria ('Combustoleo') en la tabla de Combustoleo.")

            categories_comb = []
            proceso_vals_comb = []
            diario_vals_comb = []
            programa_vals_comb = []
            columna1_vals_comb = []

            # Filtrar df_prod_combustoleo para quedarnos con los años y meses activos
            prod_rows_comb = []
            for idx, row in app.df_prod_combustoleo.iterrows():
                cat = str(row.iloc[0]).strip()
                val = row.iloc[1]

                if not cat:
                    continue

                if not any(c.isalpha() for c in cat):
                    prod_rows_comb.append((cat, val))
                else:
                    try:
                        p_val = float(val)
                    except:
                        p_val = 0
                    # Permitir meses con producción 0 para no saltar el mes actual
                    prod_rows_comb.append((cat, val))

            # Ajustar al límite de 30 categorías
            if len(prod_rows_comb) > 30:
                prod_rows_comb = prod_rows_comb[-30:]

            # Llenar filas de años y meses (sin celdas vacías al final)
            for i in range(len(prod_rows_comb)):
                cat_val = prod_rows_comb[i][0]
                try:
                    proc_val = float(prod_rows_comb[i][1])
                except:
                    proc_val = None

                categories_comb.append(cat_val)
                proceso_vals_comb.append(proc_val)
                diario_vals_comb.append(None)
                programa_vals_comb.append(None)
                columna1_vals_comb.append(None)

            # Llenar filas diarias
            for i in range(31):
                categories_comb.append(str(i + 1))
                proceso_vals_comb.append(None)

                d_val = None
                if i < len(app.df_data_combustoleo):
                    try:
                        d_val = float(app.df_data_combustoleo[diario_col_comb].iloc[i])
                    except:
                        d_val = None
                diario_vals_comb.append(d_val)

                p_val = None
                if i < len(app.df_snr_combustoleo):
                    try:
                        p_val = float(app.df_snr_combustoleo.iloc[i, 0])
                    except:
                        p_val = None
                programa_vals_comb.append(p_val)

                c_val = None
                if i < len(app.df_snr_combustoleo):
                    try:
                        c_val = float(app.df_snr_combustoleo.iloc[i, 1])
                    except:
                        c_val = None
                columna1_vals_comb.append(c_val)

            # Actualizar la gráfica de Combustoleo (Diapositiva 7)
            update_slide_chart(chart_comb, categories_comb, proceso_vals_comb, diario_vals_comb, programa_vals_comb, columna1_vals_comb, wine_color, green_color, current_month_name=master_current_month)
 
            # --- 7. PROCESAR DIAPOSITIVA DE GASOLINAS CADEREYTA (DIAPOSITIVA 11) ---
            if app.df_data_cad_gas is not None and app.df_snr_cad_gas is not None and app.df_prod_cad_gas is not None:
                slide_cad_gas = prs.slides[10]
                chart_cad_gas = None
                for shape in slide_cad_gas.shapes:
                    if shape.has_chart:
                        chart_cad_gas = shape.chart
                        break
                
                if chart_cad_gas:
                    # Buscar columna SNR o usar la segunda por defecto
                    snr_col_cad_gas = None
                    for col in app.df_data_cad_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_cad_gas = col
                            break
                    if not snr_col_cad_gas and len(app.df_data_cad_gas.columns) >= 2:
                        snr_col_cad_gas = app.df_data_cad_gas.columns[1]
                    
                    if snr_col_cad_gas:
                        categories_cg = []
                        proceso_vals_cg = []
                        diario_vals_cg = []
                        programa_vals_cg = []
                        columna1_vals_cg = []

                        # Filtrar producción anual (últimas 30 categorías)
                        prod_rows_cg = []
                        for idx, row in app.df_prod_cad_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_cg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_cg.append((cat, val))
                                except: pass





                        if len(prod_rows_cg) > 30: prod_rows_cg = prod_rows_cg[-30:]

                        for i in range(len(prod_rows_cg)):
                            categories_cg.append(prod_rows_cg[i][0])
                            try: proceso_vals_cg.append(float(prod_rows_cg[i][1]))
                            except: proceso_vals_cg.append(None)
                            diario_vals_cg.append(None)
                            programa_vals_cg.append(None)
                            columna1_vals_cg.append(None)
 
                        # Llenar datos diarios (31 días)
                        for i in range(31):
                            categories_cg.append(str(i + 1))
                            proceso_vals_cg.append(None)
                            
                            try: diario_vals_cg.append(float(app.df_data_cad_gas[snr_col_cad_gas].iloc[i]))
                            except: diario_vals_cg.append(None)
                            
                            try: programa_vals_cg.append(float(app.df_snr_cad_gas.iloc[i, 0]))
                            except: programa_vals_cg.append(None)
                            
                            try: columna1_vals_cg.append(float(app.df_snr_cad_gas.iloc[i, 1]))
                            except: columna1_vals_cg.append(None)
 
                        update_slide_chart(chart_cad_gas, categories_cg, proceso_vals_cg, diario_vals_cg, programa_vals_cg, columna1_vals_cg, wine_color, green_color)
 
                # --- 8. PROCESAR DIAPOSITIVA DE DIESEL CADEREYTA (DIAPOSITIVA 12) ---
                if app.df_data_cad_die is not None and app.df_snr_cad_die is not None and app.df_prod_cad_die is not None:
                    slide_cad_die = prs.slides[11]
                    chart_cad_die = None
                    for shape in slide_cad_die.shapes:
                        if shape.has_chart:
                            chart_cad_die = shape.chart
                            break
                    
                    if chart_cad_die:
                        snr_col_cad_die = None
                        for col in app.df_data_cad_die.columns:
                            if "SNR" in str(col).upper():
                                snr_col_cad_die = col
                                break
                        if not snr_col_cad_die and len(app.df_data_cad_die.columns) >= 2:
                            snr_col_cad_die = app.df_data_cad_die.columns[1]
                        
                        if snr_col_cad_die:
                            categories_cd = []
                            proceso_vals_cd = []
                            diario_vals_cd = []
                            programa_vals_cd = []
                            columna1_vals_cd = []
 
                            prod_rows_cd = []
                            for idx, row in app.df_prod_cad_die.iterrows():
                                cat = str(row.iloc[0]).strip()
                                val = row.iloc[1]
                                if not cat: continue
                                if not any(c.isalpha() for c in cat):
                                    prod_rows_cd.append((cat, val))
                                else:
                                    try:
                                        if float(val) != 0: prod_rows_cd.append((cat, val))
                                    except: pass
                            
                            if len(prod_rows_cd) > 30: prod_rows_cd = prod_rows_cd[-30:]
 
                            for i in range(len(prod_rows_cd)):
                                categories_cd.append(prod_rows_cd[i][0])
                                try: proceso_vals_cd.append(float(prod_rows_cd[i][1]))
                                except: proceso_vals_cd.append(None)
                                diario_vals_cd.append(None)
                                programa_vals_cd.append(None)
                                columna1_vals_cd.append(None)
 
                            # Llenar datos diarios (31 días)
                            for i in range(31):
                                categories_cd.append(str(i + 1))
                                proceso_vals_cd.append(None)
                                
                                try: diario_vals_cd.append(float(app.df_data_cad_die[snr_col_cad_die].iloc[i]))
                                except: diario_vals_cd.append(None)
                                
                                try: programa_vals_cd.append(float(app.df_snr_cad_die.iloc[i, 0]))
                                except: programa_vals_cd.append(None)
                                
                                try: columna1_vals_cd.append(float(app.df_snr_cad_die.iloc[i, 1]))
                                except: columna1_vals_cd.append(None)
 
                            update_slide_chart(chart_cad_die, categories_cd, proceso_vals_cd, diario_vals_cd, programa_vals_cd, columna1_vals_cd, wine_color, green_color)

                # --- 9. PROCESAR DIAPOSITIVA DE COMBUSTOLEO CADEREYTA (DIAPOSITIVA 13) ---
                if len(prs.slides) > 12 and app.df_data_cad_comb is not None and app.df_snr_cad_comb is not None and app.df_prod_cad_comb is not None:
                    slide_cad_comb = prs.slides[12]
                    chart_cad_comb = None
                    for shape in slide_cad_comb.shapes:
                        if shape.has_chart:
                            chart_cad_comb = shape.chart
                            break
                    
                    if chart_cad_comb:
                        snr_col_cad_comb = None
                        for col in app.df_data_cad_comb.columns:
                            if "SNR" in str(col).upper():
                                snr_col_cad_comb = col
                                break
                        if not snr_col_cad_comb and len(app.df_data_cad_comb.columns) >= 2:
                            snr_col_cad_comb = app.df_data_cad_comb.columns[1]
                        
                        if snr_col_cad_comb:
                            categories_cc = []
                            proceso_vals_cc = []
                            diario_vals_cc = []
                            programa_vals_cc = []
                            columna1_vals_cc = []
 
                            prod_rows_cc = []
                            for idx, row in app.df_prod_cad_comb.iterrows():
                                cat = str(row.iloc[0]).strip()
                                val = row.iloc[1]
                                if not cat: continue
                                if not any(c.isalpha() for c in cat):
                                    prod_rows_cc.append((cat, val))
                                else:
                                    try:
                                        if float(val) != 0: prod_rows_cc.append((cat, val))
                                    except: pass
                            
                            if len(prod_rows_cc) > 30: prod_rows_cc = prod_rows_cc[-30:]
 
                            for i in range(len(prod_rows_cc)):
                                categories_cc.append(prod_rows_cc[i][0])
                                try: proceso_vals_cc.append(float(prod_rows_cc[i][1]))
                                except: proceso_vals_cc.append(None)
                                diario_vals_cc.append(None)
                                programa_vals_cc.append(None)
                                columna1_vals_cc.append(None)
 
                            # Llenar datos diarios (31 días)
                            for i in range(31):
                                categories_cc.append(str(i + 1))
                                proceso_vals_cc.append(None)
                                
                                try: diario_vals_cc.append(float(app.df_data_cad_comb[snr_col_cad_comb].iloc[i]))
                                except: diario_vals_cc.append(None)
                                
                                try: programa_vals_cc.append(float(app.df_snr_cad_comb.iloc[i, 0]))
                                except: programa_vals_cc.append(None)
                                
                                try: columna1_vals_cc.append(float(app.df_snr_cad_comb.iloc[i, 1]))
                                except: columna1_vals_cc.append(None)
 
                            update_slide_chart(chart_cad_comb, categories_cc, proceso_vals_cc, diario_vals_cc, programa_vals_cc, columna1_vals_cc, wine_color, green_color)

            # --- 10. PROCESAR DIAPOSITIVA DE CRUDO MADERO (DIAPOSITIVA 16) ---
            if len(prs.slides) > 15 and app.df_data_mad_crud is not None and app.df_snr_mad_crud is not None and app.df_prod_mad_crud is not None:
                slide_mad_crud = prs.slides[15]
                chart_mad_crud = None
                for shape in slide_mad_crud.shapes:
                    if shape.has_chart:
                        chart_mad_crud = shape.chart
                        break
                
                if chart_mad_crud:
                    categories_mc = []
                    proceso_vals_mc = []
                    diario_vals_mc = []
                    programa_vals_mc = []
                    columna1_vals_mc = []

                    prod_rows_mc = []
                    for idx, row in app.df_prod_mad_crud.iterrows():
                        cat = str(row.iloc[0]).strip()
                        val = row.iloc[1]
                        if not cat: continue
                        if not any(c.isalpha() for c in cat):
                            prod_rows_mc.append((cat, val))
                        else:
                            try:
                                if float(val) != 0: prod_rows_mc.append((cat, val))
                            except: pass

                    if len(prod_rows_mc) > 30: prod_rows_mc = prod_rows_mc[-30:]

                    for i in range(len(prod_rows_mc)):
                        categories_mc.append(prod_rows_mc[i][0])
                        try: proceso_vals_mc.append(float(prod_rows_mc[i][1]))
                        except: proceso_vals_mc.append(None)
                        diario_vals_mc.append(None)
                        programa_vals_mc.append(None)
                        columna1_vals_mc.append(None)

                    proceso_col_mad_crud = app.df_data_mad_crud.columns[1]

                    for i in range(31):
                        categories_mc.append(str(i + 1))
                        proceso_vals_mc.append(None)
                        
                        try: diario_vals_mc.append(float(app.df_data_mad_crud[proceso_col_mad_crud].iloc[i]))
                        except: diario_vals_mc.append(None)
                        
                        try: programa_vals_mc.append(float(app.df_snr_mad_crud.iloc[i, 0]))
                        except: programa_vals_mc.append(None)
                        
                        try: columna1_vals_mc.append(float(app.df_snr_mad_crud.iloc[i, 1]))
                        except: columna1_vals_mc.append(None)

                    update_slide_chart(chart_mad_crud, categories_mc, proceso_vals_mc, diario_vals_mc, programa_vals_mc, columna1_vals_mc, wine_color, green_color)

            # --- 11. PROCESAR DIAPOSITIVA DE GASOLINAS MADERO (DIAPOSITIVA 17) ---
            if len(prs.slides) > 16 and app.df_data_mad_gas is not None and app.df_snr_mad_gas is not None and app.df_prod_mad_gas is not None:
                slide_mad_gas = prs.slides[16]
                chart_mad_gas = None
                for shape in slide_mad_gas.shapes:
                    if shape.has_chart:
                        chart_mad_gas = shape.chart
                        break
                
                if chart_mad_gas:
                    snr_col_mad_gas = None
                    for col in app.df_data_mad_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mad_gas = col
                            break
                    if not snr_col_mad_gas and len(app.df_data_mad_gas.columns) >= 2:
                        snr_col_mad_gas = app.df_data_mad_gas.columns[1]

                    if snr_col_mad_gas:
                        categories_mg = []
                        proceso_vals_mg = []
                        diario_vals_mg = []
                        programa_vals_mg = []
                        columna1_vals_mg = []

                        prod_rows_mg = []
                        for idx, row in app.df_prod_mad_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mg.append((cat, val))
                                except: pass

                        if len(prod_rows_mg) > 30: prod_rows_mg = prod_rows_mg[-30:]

                        for i in range(len(prod_rows_mg)):
                            categories_mg.append(prod_rows_mg[i][0])
                            try: proceso_vals_mg.append(float(prod_rows_mg[i][1]))
                            except: proceso_vals_mg.append(None)
                            diario_vals_mg.append(None)
                            programa_vals_mg.append(None)
                            columna1_vals_mg.append(None)

                        for i in range(31):
                            categories_mg.append(str(i + 1))
                            proceso_vals_mg.append(None)
                            
                            try: diario_vals_mg.append(float(app.df_data_mad_gas[snr_col_mad_gas].iloc[i]))
                            except: diario_vals_mg.append(None)
                            
                            try: programa_vals_mg.append(float(app.df_snr_mad_gas.iloc[i, 0]))
                            except: programa_vals_mg.append(None)
                            
                            try: columna1_vals_mg.append(float(app.df_snr_mad_gas.iloc[i, 1]))
                            except: columna1_vals_mg.append(None)

                        update_slide_chart(chart_mad_gas, categories_mg, proceso_vals_mg, diario_vals_mg, programa_vals_mg, columna1_vals_mg, wine_color, green_color)

            # --- 12. PROCESAR DIAPOSITIVA DE DIESEL MADERO (DIAPOSITIVA 18) ---
            if len(prs.slides) > 17 and app.df_data_mad_die is not None and app.df_snr_mad_die is not None and app.df_prod_mad_die is not None:
                slide_mad_die = prs.slides[17]
                chart_mad_die = None
                for shape in slide_mad_die.shapes:
                    if shape.has_chart:
                        chart_mad_die = shape.chart
                        break
                
                if chart_mad_die:
                    snr_col_mad_die = None
                    for col in app.df_data_mad_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mad_die = col
                            break
                    if not snr_col_mad_die and len(app.df_data_mad_die.columns) >= 2:
                        snr_col_mad_die = app.df_data_mad_die.columns[1]

                    if snr_col_mad_die:
                        categories_md = []
                        proceso_vals_md = []
                        diario_vals_md = []
                        programa_vals_md = []
                        columna1_vals_md = []

                        prod_rows_md = []
                        for idx, row in app.df_prod_mad_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_md.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_md.append((cat, val))
                                except: pass

                        if len(prod_rows_md) > 30: prod_rows_md = prod_rows_md[-30:]

                        for i in range(len(prod_rows_md)):
                            categories_md.append(prod_rows_md[i][0])
                            try: proceso_vals_md.append(float(prod_rows_md[i][1]))
                            except: proceso_vals_md.append(None)
                            diario_vals_md.append(None)
                            programa_vals_md.append(None)
                            columna1_vals_md.append(None)

                        for i in range(31):
                            categories_md.append(str(i + 1))
                            proceso_vals_md.append(None)
                            
                            try: diario_vals_md.append(float(app.df_data_mad_die[snr_col_mad_die].iloc[i]))
                            except: diario_vals_md.append(None)
                            
                            try: programa_vals_md.append(float(app.df_snr_mad_die.iloc[i, 0]))
                            except: programa_vals_md.append(None)
                            
                            try: columna1_vals_md.append(float(app.df_snr_mad_die.iloc[i, 1]))
                            except: columna1_vals_md.append(None)

                        update_slide_chart(chart_mad_die, categories_md, proceso_vals_md, diario_vals_md, programa_vals_md, columna1_vals_md, wine_color, green_color)

            # --- 13. PROCESAR DIAPOSITIVA DE TURBOSINA MADERO (DIAPOSITIVA 19) ---
            if len(prs.slides) > 18 and app.df_data_mad_turb is not None and app.df_snr_mad_turb is not None and app.df_prod_mad_turb is not None:
                slide_mad_turb = prs.slides[18]
                chart_mad_turb = None
                for shape in slide_mad_turb.shapes:
                    if shape.has_chart:
                        chart_mad_turb = shape.chart
                        break
                
                if chart_mad_turb:
                    snr_col_mad_turb = None
                    for col in app.df_data_mad_turb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mad_turb = col
                            break
                    if not snr_col_mad_turb and len(app.df_data_mad_turb.columns) >= 2:
                        snr_col_mad_turb = app.df_data_mad_turb.columns[1]

                    if snr_col_mad_turb:
                        categories_mtu = []
                        proceso_vals_mtu = []
                        diario_vals_mtu = []
                        programa_vals_mtu = []
                        columna1_vals_mtu = []

                        prod_rows_mtu = []
                        for idx, row in app.df_prod_mad_turb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mtu.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mtu.append((cat, val))
                                except: pass

                        if len(prod_rows_mtu) > 30: prod_rows_mtu = prod_rows_mtu[-30:]

                        for i in range(len(prod_rows_mtu)):
                            categories_mtu.append(prod_rows_mtu[i][0])
                            try: proceso_vals_mtu.append(float(prod_rows_mtu[i][1]))
                            except: proceso_vals_mtu.append(None)
                            diario_vals_mtu.append(None)
                            programa_vals_mtu.append(None)
                            columna1_vals_mtu.append(None)

                        for i in range(31):
                            categories_mtu.append(str(i + 1))
                            proceso_vals_mtu.append(None)
                            
                            try: diario_vals_mtu.append(float(app.df_data_mad_turb[snr_col_mad_turb].iloc[i]))
                            except: diario_vals_mtu.append(None)
                            
                            try: programa_vals_mtu.append(float(app.df_snr_mad_turb.iloc[i, 0]))
                            except: programa_vals_mtu.append(None)
                            
                            try: columna1_vals_mtu.append(float(app.df_snr_mad_turb.iloc[i, 1]))
                            except: columna1_vals_mtu.append(None)

                        update_slide_chart(chart_mad_turb, categories_mtu, proceso_vals_mtu, diario_vals_mtu, programa_vals_mtu, columna1_vals_mtu, wine_color, green_color)

            # --- 14. PROCESAR DIAPOSITIVA DE COMBUSTOLEO MADERO (DIAPOSITIVA 20) ---
            if len(prs.slides) > 19 and app.df_data_mad_comb is not None and app.df_snr_mad_comb is not None and app.df_prod_mad_comb is not None:
                slide_mad_comb = prs.slides[19]
                chart_mad_comb = None
                for shape in slide_mad_comb.shapes:
                    if shape.has_chart:
                        chart_mad_comb = shape.chart
                        break
                
                if chart_mad_comb:
                    snr_col_mad_comb = None
                    for col in app.df_data_mad_comb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mad_comb = col
                            break
                    if not snr_col_mad_comb and len(app.df_data_mad_comb.columns) >= 2:
                        snr_col_mad_comb = app.df_data_mad_comb.columns[1]

                    if snr_col_mad_comb:
                        categories_mco = []
                        proceso_vals_mco = []
                        diario_vals_mco = []
                        programa_vals_mco = []
                        columna1_vals_mco = []

                        prod_rows_mco = []
                        for idx, row in app.df_prod_mad_comb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mco.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mco.append((cat, val))
                                except: pass

                        if len(prod_rows_mco) > 30: prod_rows_mco = prod_rows_mco[-30:]

                        for i in range(len(prod_rows_mco)):
                            categories_mco.append(prod_rows_mco[i][0])
                            try: proceso_vals_mco.append(float(prod_rows_mco[i][1]))
                            except: proceso_vals_mco.append(None)
                            diario_vals_mco.append(None)
                            programa_vals_mco.append(None)
                            columna1_vals_mco.append(None)

                        for i in range(31):
                            categories_mco.append(str(i + 1))
                            proceso_vals_mco.append(None)
                            
                            try: diario_vals_mco.append(float(app.df_data_mad_comb[snr_col_mad_comb].iloc[i]))
                            except: diario_vals_mco.append(None)
                            
                            try: programa_vals_mco.append(float(app.df_snr_mad_comb.iloc[i, 0]))
                            except: programa_vals_mco.append(None)
                            
                            try: columna1_vals_mco.append(float(app.df_snr_mad_comb.iloc[i, 1]))
                            except: columna1_vals_mco.append(None)

                        update_slide_chart(chart_mad_comb, categories_mco, proceso_vals_mco, diario_vals_mco, programa_vals_mco, columna1_vals_mco, wine_color, green_color)

            # --- 15. PROCESAR DIAPOSITIVA DE CRUDO MINATITLAN (DIAPOSITIVA 23) ---
            if len(prs.slides) > 22 and app.df_data_mina_crud is not None and app.df_snr_mina_crud is not None and app.df_prod_mina_crud is not None:
                slide_mina_crud = prs.slides[22]
                chart_mina_crud = None
                for shape in slide_mina_crud.shapes:
                    if shape.has_chart:
                        chart_mina_crud = shape.chart
                        break
                
                if chart_mina_crud:
                    snr_col_mina_crud = None
                    for col in app.df_data_mina_crud.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mina_crud = col
                            break
                    if not snr_col_mina_crud and len(app.df_data_mina_crud.columns) >= 2:
                        snr_col_mina_crud = app.df_data_mina_crud.columns[1]

                    if snr_col_mina_crud:
                        categories_mic = []
                        proceso_vals_mic = []
                        diario_vals_mic = []
                        programa_vals_mic = []
                        columna1_vals_mic = []

                        prod_rows_mic = []
                        for idx, row in app.df_prod_mina_crud.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mic.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mic.append((cat, val))
                                except: pass

                        if len(prod_rows_mic) > 30: prod_rows_mic = prod_rows_mic[-30:]

                        for i in range(len(prod_rows_mic)):
                            categories_mic.append(prod_rows_mic[i][0])
                            try: proceso_vals_mic.append(float(prod_rows_mic[i][1]))
                            except: proceso_vals_mic.append(None)
                            diario_vals_mic.append(None)
                            programa_vals_mic.append(None)
                            columna1_vals_mic.append(None)

                        for i in range(31):
                            categories_mic.append(str(i + 1))
                            proceso_vals_mic.append(None)
                            
                            try: diario_vals_mic.append(float(app.df_data_mina_crud[snr_col_mina_crud].iloc[i]))
                            except: diario_vals_mic.append(None)
                            
                            try: programa_vals_mic.append(float(app.df_snr_mina_crud.iloc[i, 0]))
                            except: programa_vals_mic.append(None)
                            
                            try: columna1_vals_mic.append(float(app.df_snr_mina_crud.iloc[i, 1]))
                            except: columna1_vals_mic.append(None)

                        update_slide_chart(chart_mina_crud, categories_mic, proceso_vals_mic, diario_vals_mic, programa_vals_mic, columna1_vals_mic, wine_color, green_color)

            # --- 16. PROCESAR DIAPOSITIVA DE GASOLINAS MINATITLAN (DIAPOSITIVA 24) ---
            if len(prs.slides) > 23 and app.df_data_mina_gas is not None and app.df_snr_mina_gas is not None and app.df_prod_mina_gas is not None:
                slide_mina_gas = prs.slides[23]
                chart_mina_gas = None
                for shape in slide_mina_gas.shapes:
                    if shape.has_chart:
                        chart_mina_gas = shape.chart
                        break
                
                if chart_mina_gas:
                    snr_col_mina_gas = None
                    for col in app.df_data_mina_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mina_gas = col
                            break
                    if not snr_col_mina_gas and len(app.df_data_mina_gas.columns) >= 2:
                        snr_col_mina_gas = app.df_data_mina_gas.columns[1]

                    if snr_col_mina_gas:
                        categories_mig = []
                        proceso_vals_mig = []
                        diario_vals_mig = []
                        programa_vals_mig = []
                        columna1_vals_mig = []

                        prod_rows_mig = []
                        for idx, row in app.df_prod_mina_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mig.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mig.append((cat, val))
                                except: pass

                        if len(prod_rows_mig) > 30: prod_rows_mig = prod_rows_mig[-30:]

                        for i in range(len(prod_rows_mig)):
                            categories_mig.append(prod_rows_mig[i][0])
                            try: proceso_vals_mig.append(float(prod_rows_mig[i][1]))
                            except: proceso_vals_mig.append(None)
                            diario_vals_mig.append(None)
                            programa_vals_mig.append(None)
                            columna1_vals_mig.append(None)

                        for i in range(31):
                            categories_mig.append(str(i + 1))
                            proceso_vals_mig.append(None)
                            
                            try: diario_vals_mig.append(float(app.df_data_mina_gas[snr_col_mina_gas].iloc[i]))
                            except: diario_vals_mig.append(None)
                            
                            try: programa_vals_mig.append(float(app.df_snr_mina_gas.iloc[i, 0]))
                            except: programa_vals_mig.append(None)
                            
                            try: columna1_vals_mig.append(float(app.df_snr_mina_gas.iloc[i, 1]))
                            except: columna1_vals_mig.append(None)

                        update_slide_chart(chart_mina_gas, categories_mig, proceso_vals_mig, diario_vals_mig, programa_vals_mig, columna1_vals_mig, wine_color, green_color)

            # --- 17. PROCESAR DIAPOSITIVA DE DIESEL MINATITLAN (DIAPOSITIVA 25) ---
            if len(prs.slides) > 24 and app.df_data_mina_die is not None and app.df_snr_mina_die is not None and app.df_prod_mina_die is not None:
                slide_mina_die = prs.slides[24]
                chart_mina_die = None
                for shape in slide_mina_die.shapes:
                    if shape.has_chart:
                        chart_mina_die = shape.chart
                        break
                
                if chart_mina_die:
                    snr_col_mina_die = None
                    for col in app.df_data_mina_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mina_die = col
                            break
                    if not snr_col_mina_die and len(app.df_data_mina_die.columns) >= 2:
                        snr_col_mina_die = app.df_data_mina_die.columns[1]

                    if snr_col_mina_die:
                        categories_mid = []
                        proceso_vals_mid = []
                        diario_vals_mid = []
                        programa_vals_mid = []
                        columna1_vals_mid = []

                        prod_rows_mid = []
                        for idx, row in app.df_prod_mina_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mid.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mid.append((cat, val))
                                except: pass

                        if len(prod_rows_mid) > 30: prod_rows_mid = prod_rows_mid[-30:]

                        for i in range(len(prod_rows_mid)):
                            categories_mid.append(prod_rows_mid[i][0])
                            try: proceso_vals_mid.append(float(prod_rows_mid[i][1]))
                            except: proceso_vals_mid.append(None)
                            diario_vals_mid.append(None)
                            programa_vals_mid.append(None)
                            columna1_vals_mid.append(None)

                        for i in range(31):
                            categories_mid.append(str(i + 1))
                            proceso_vals_mid.append(None)
                            
                            try: diario_vals_mid.append(float(app.df_data_mina_die[snr_col_mina_die].iloc[i]))
                            except: diario_vals_mid.append(None)
                            
                            try: programa_vals_mid.append(float(app.df_snr_mina_die.iloc[i, 0]))
                            except: programa_vals_mid.append(None)
                            
                            try: columna1_vals_mid.append(float(app.df_snr_mina_die.iloc[i, 1]))
                            except: columna1_vals_mid.append(None)

                        update_slide_chart(chart_mina_die, categories_mid, proceso_vals_mid, diario_vals_mid, programa_vals_mid, columna1_vals_mid, wine_color, green_color)

            # --- 18. PROCESAR DIAPOSITIVA DE COMBUSTOLEO MINATITLAN (DIAPOSITIVA 26) ---
            if len(prs.slides) > 25 and app.df_data_mina_comb is not None and app.df_snr_mina_comb is not None and app.df_prod_mina_comb is not None:
                slide_mina_comb = prs.slides[25]
                chart_mina_comb = None
                for shape in slide_mina_comb.shapes:
                    if shape.has_chart:
                        chart_mina_comb = shape.chart
                        break
                
                if chart_mina_comb:
                    snr_col_mina_comb = None
                    for col in app.df_data_mina_comb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_mina_comb = col
                            break
                    if not snr_col_mina_comb and len(app.df_data_mina_comb.columns) >= 2:
                        snr_col_mina_comb = app.df_data_mina_comb.columns[1]

                    if snr_col_mina_comb:
                        categories_mco = []
                        proceso_vals_mco = []
                        diario_vals_mco = []
                        programa_vals_mco = []
                        columna1_vals_mco = []

                        prod_rows_mco = []
                        for idx, row in app.df_prod_mina_comb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_mco.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_mco.append((cat, val))
                                except: pass

                        if len(prod_rows_mco) > 30: prod_rows_mco = prod_rows_mco[-30:]

                        for i in range(len(prod_rows_mco)):
                            categories_mco.append(prod_rows_mco[i][0])
                            try: proceso_vals_mco.append(float(prod_rows_mco[i][1]))
                            except: proceso_vals_mco.append(None)
                            diario_vals_mco.append(None)
                            programa_vals_mco.append(None)
                            columna1_vals_mco.append(None)

                        for i in range(31):
                            categories_mco.append(str(i + 1))
                            proceso_vals_mco.append(None)
                            
                            try: diario_vals_mco.append(float(app.df_data_mina_comb[snr_col_mina_comb].iloc[i]))
                            except: diario_vals_mco.append(None)
                            
                            try: programa_vals_mco.append(float(app.df_snr_mina_comb.iloc[i, 0]))
                            except: programa_vals_mco.append(None)
                            
                            try: columna1_vals_mco.append(float(app.df_snr_mina_comb.iloc[i, 1]))
                            except: columna1_vals_mco.append(None)

                        update_slide_chart(chart_mina_comb, categories_mco, proceso_vals_mco, diario_vals_mco, programa_vals_mco, columna1_vals_mco, wine_color, green_color)

            # --- 19. PROCESAR DIAPOSITIVA DE CRUDO SALAMANCA (DIAPOSITIVA 29) ---
            if len(prs.slides) > 28 and app.df_data_sala_crud is not None and app.df_snr_sala_crud is not None and app.df_prod_sala_crud is not None:
                slide_sala_crud = prs.slides[28]
                chart_sala_crud = None
                for shape in slide_sala_crud.shapes:
                    if shape.has_chart:
                        chart_sala_crud = shape.chart
                        break
                
                if chart_sala_crud:
                    snr_col_sala_crud = None
                    for col in app.df_data_sala_crud.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sala_crud = col
                            break
                    if not snr_col_sala_crud and len(app.df_data_sala_crud.columns) >= 2:
                        snr_col_sala_crud = app.df_data_sala_crud.columns[1]

                    if snr_col_sala_crud:
                        categories_sc = []
                        proceso_vals_sc = []
                        diario_vals_sc = []
                        programa_vals_sc = []
                        columna1_vals_sc = []

                        prod_rows_sc = []
                        for idx, row in app.df_prod_sala_crud.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sc.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sc.append((cat, val))
                                except: pass

                        if len(prod_rows_sc) > 30: prod_rows_sc = prod_rows_sc[-30:]

                        for i in range(len(prod_rows_sc)):
                            categories_sc.append(prod_rows_sc[i][0])
                            try: proceso_vals_sc.append(float(prod_rows_sc[i][1]))
                            except: proceso_vals_sc.append(None)
                            diario_vals_sc.append(None)
                            programa_vals_sc.append(None)
                            columna1_vals_sc.append(None)

                        for i in range(31):
                            categories_sc.append(str(i + 1))
                            proceso_vals_sc.append(None)
                            
                            try: diario_vals_sc.append(float(app.df_data_sala_crud[snr_col_sala_crud].iloc[i]))
                            except: diario_vals_sc.append(None)
                            
                            try: programa_vals_sc.append(float(app.df_snr_sala_crud.iloc[i, 0]))
                            except: programa_vals_sc.append(None)
                            
                            try: columna1_vals_sc.append(float(app.df_snr_sala_crud.iloc[i, 1]))
                            except: columna1_vals_sc.append(None)

                        update_slide_chart(chart_sala_crud, categories_sc, proceso_vals_sc, diario_vals_sc, programa_vals_sc, columna1_vals_sc, wine_color, green_color)

            # --- 20. PROCESAR DIAPOSITIVA DE GASOLINAS SALAMANCA (DIAPOSITIVA 30) ---
            if len(prs.slides) > 29 and app.df_data_sala_gas is not None and app.df_snr_sala_gas is not None and app.df_prod_sala_gas is not None:
                slide_sala_gas = prs.slides[29]
                chart_sala_gas = None
                for shape in slide_sala_gas.shapes:
                    if shape.has_chart:
                        chart_sala_gas = shape.chart
                        break
                
                if chart_sala_gas:
                    snr_col_sala_gas = None
                    for col in app.df_data_sala_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sala_gas = col
                            break
                    if not snr_col_sala_gas and len(app.df_data_sala_gas.columns) >= 2:
                        snr_col_sala_gas = app.df_data_sala_gas.columns[1]

                    if snr_col_sala_gas:
                        categories_sg = []
                        proceso_vals_sg = []
                        diario_vals_sg = []
                        programa_vals_sg = []
                        columna1_vals_sg = []

                        prod_rows_sg = []
                        for idx, row in app.df_prod_sala_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sg.append((cat, val))
                                except: pass

                        if len(prod_rows_sg) > 30: prod_rows_sg = prod_rows_sg[-30:]

                        for i in range(len(prod_rows_sg)):
                            categories_sg.append(prod_rows_sg[i][0])
                            try: proceso_vals_sg.append(float(prod_rows_sg[i][1]))
                            except: proceso_vals_sg.append(None)
                            diario_vals_sg.append(None)
                            programa_vals_sg.append(None)
                            columna1_vals_sg.append(None)

                        for i in range(31):
                            categories_sg.append(str(i + 1))
                            proceso_vals_sg.append(None)
                            
                            try: diario_vals_sg.append(float(app.df_data_sala_gas[snr_col_sala_gas].iloc[i]))
                            except: diario_vals_sg.append(None)
                            
                            try: programa_vals_sg.append(float(app.df_snr_sala_gas.iloc[i, 0]))
                            except: programa_vals_sg.append(None)
                            
                            try: columna1_vals_sg.append(float(app.df_snr_sala_gas.iloc[i, 1]))
                            except: columna1_vals_sg.append(None)

                        update_slide_chart(chart_sala_gas, categories_sg, proceso_vals_sg, diario_vals_sg, programa_vals_sg, columna1_vals_sg, wine_color, green_color)

            # --- 21. PROCESAR DIAPOSITIVA DE DIESEL SALAMANCA (DIAPOSITIVA 31) ---
            if len(prs.slides) > 30 and app.df_data_sala_die is not None and app.df_snr_sala_die is not None and app.df_prod_sala_die is not None:
                slide_sala_die = prs.slides[30]
                chart_sala_die = None
                for shape in slide_sala_die.shapes:
                    if shape.has_chart:
                        chart_sala_die = shape.chart
                        break
                
                if chart_sala_die:
                    snr_col_sala_die = None
                    for col in app.df_data_sala_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sala_die = col
                            break
                    if not snr_col_sala_die and len(app.df_data_sala_die.columns) >= 2:
                        snr_col_sala_die = app.df_data_sala_die.columns[1]

                    if snr_col_sala_die:
                        categories_sd = []
                        proceso_vals_sd = []
                        diario_vals_sd = []
                        programa_vals_sd = []
                        columna1_vals_sd = []

                        prod_rows_sd = []
                        for idx, row in app.df_prod_sala_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sd.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sd.append((cat, val))
                                except: pass

                        if len(prod_rows_sd) > 30: prod_rows_sd = prod_rows_sd[-30:]

                        for i in range(len(prod_rows_sd)):
                            categories_sd.append(prod_rows_sd[i][0])
                            try: proceso_vals_sd.append(float(prod_rows_sd[i][1]))
                            except: proceso_vals_sd.append(None)
                            diario_vals_sd.append(None)
                            programa_vals_sd.append(None)
                            columna1_vals_sd.append(None)

                        for i in range(31):
                            categories_sd.append(str(i + 1))
                            proceso_vals_sd.append(None)
                            
                            try: diario_vals_sd.append(float(app.df_data_sala_die[snr_col_sala_die].iloc[i]))
                            except: diario_vals_sd.append(None)
                            
                            try: programa_vals_sd.append(float(app.df_snr_sala_die.iloc[i, 0]))
                            except: programa_vals_sd.append(None)
                            
                            try: columna1_vals_sd.append(float(app.df_snr_sala_die.iloc[i, 1]))
                            except: columna1_vals_sd.append(None)

                        update_slide_chart(chart_sala_die, categories_sd, proceso_vals_sd, diario_vals_sd, programa_vals_sd, columna1_vals_sd, wine_color, green_color)

            # --- 22. PROCESAR DIAPOSITIVA DE TURBOSINA SALAMANCA (DIAPOSITIVA 32) ---
            if len(prs.slides) > 31 and app.df_data_sala_turb is not None and app.df_snr_sala_turb is not None and app.df_prod_sala_turb is not None:
                slide_sala_turb = prs.slides[31]
                chart_sala_turb = None
                for shape in slide_sala_turb.shapes:
                    if shape.has_chart:
                        chart_sala_turb = shape.chart
                        break
                
                if chart_sala_turb:
                    snr_col_sala_turb = None
                    for col in app.df_data_sala_turb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sala_turb = col
                            break
                    if not snr_col_sala_turb and len(app.df_data_sala_turb.columns) >= 2:
                        snr_col_sala_turb = app.df_data_sala_turb.columns[1]

                    if snr_col_sala_turb:
                        categories_stur = []
                        proceso_vals_stur = []
                        diario_vals_stur = []
                        programa_vals_stur = []
                        columna1_vals_stur = []

                        prod_rows_stur = []
                        for idx, row in app.df_prod_sala_turb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_stur.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_stur.append((cat, val))
                                except: pass

                        if len(prod_rows_stur) > 30: prod_rows_stur = prod_rows_stur[-30:]

                        for i in range(len(prod_rows_stur)):
                            categories_stur.append(prod_rows_stur[i][0])
                            try: proceso_vals_stur.append(float(prod_rows_stur[i][1]))
                            except: proceso_vals_stur.append(None)
                            diario_vals_stur.append(None)
                            programa_vals_stur.append(None)
                            columna1_vals_stur.append(None)

                        for i in range(31):
                            categories_stur.append(str(i + 1))
                            proceso_vals_stur.append(None)
                            
                            try: diario_vals_stur.append(float(app.df_data_sala_turb[snr_col_sala_turb].iloc[i]))
                            except: diario_vals_stur.append(None)
                            
                            try: programa_vals_stur.append(float(app.df_snr_sala_turb.iloc[i, 0]))
                            except: programa_vals_stur.append(None)
                            
                            try: columna1_vals_stur.append(float(app.df_snr_sala_turb.iloc[i, 1]))
                            except: columna1_vals_stur.append(None)

                        update_slide_chart(chart_sala_turb, categories_stur, proceso_vals_stur, diario_vals_stur, programa_vals_stur, columna1_vals_stur, wine_color, green_color)

            # --- 23. PROCESAR DIAPOSITIVA DE COMBUSTOLEO SALAMANCA (DIAPOSITIVA 33) ---
            if len(prs.slides) > 32 and app.df_data_sala_comb is not None and app.df_snr_sala_comb is not None and app.df_prod_sala_comb is not None:
                slide_sala_comb = prs.slides[32]
                chart_sala_comb = None
                for shape in slide_sala_comb.shapes:
                    if shape.has_chart:
                        chart_sala_comb = shape.chart
                        break
                
                if chart_sala_comb:
                    snr_col_sala_comb = None
                    for col in app.df_data_sala_comb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sala_comb = col
                            break
                    if not snr_col_sala_comb and len(app.df_data_sala_comb.columns) >= 2:
                        snr_col_sala_comb = app.df_data_sala_comb.columns[1]

                    if snr_col_sala_comb:
                        categories_sco = []
                        proceso_vals_sco = []
                        diario_vals_sco = []
                        programa_vals_sco = []
                        columna1_vals_sco = []

                        prod_rows_sco = []
                        for idx, row in app.df_prod_sala_comb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sco.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sco.append((cat, val))
                                except: pass

                        if len(prod_rows_sco) > 30: prod_rows_sco = prod_rows_sco[-30:]

                        for i in range(len(prod_rows_sco)):
                            categories_sco.append(prod_rows_sco[i][0])
                            try: proceso_vals_sco.append(float(prod_rows_sco[i][1]))
                            except: proceso_vals_sco.append(None)
                            diario_vals_sco.append(None)
                            programa_vals_sco.append(None)
                            columna1_vals_sco.append(None)

                        for i in range(31):
                            categories_sco.append(str(i + 1))
                            proceso_vals_sco.append(None)
                            
                            try: diario_vals_sco.append(float(app.df_data_sala_comb[snr_col_sala_comb].iloc[i]))
                            except: diario_vals_sco.append(None)
                            
                            try: programa_vals_sco.append(float(app.df_snr_sala_comb.iloc[i, 0]))
                            except: programa_vals_sco.append(None)
                            
                            try: columna1_vals_sco.append(float(app.df_snr_sala_comb.iloc[i, 1]))
                            except: columna1_vals_sco.append(None)

                        update_slide_chart(chart_sala_comb, categories_sco, proceso_vals_sco, diario_vals_sco, programa_vals_sco, columna1_vals_sco, wine_color, green_color)

            # --- 24. PROCESAR DIAPOSITIVA DE CRUDO SALINA CRUZ (DIAPOSITIVA 36) ---
            if len(prs.slides) > 35 and app.df_data_sal_crud is not None and app.df_snr_sal_crud is not None and app.df_prod_sal_crud is not None:
                slide_sal_crud = prs.slides[35]
                chart_sal_crud = None
                for shape in slide_sal_crud.shapes:
                    if shape.has_chart:
                        chart_sal_crud = shape.chart
                        break
                
                if chart_sal_crud:
                    snr_col_sal_crud = None
                    for col in app.df_data_sal_crud.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sal_crud = col
                            break
                    if not snr_col_sal_crud and len(app.df_data_sal_crud.columns) >= 2:
                        snr_col_sal_crud = app.df_data_sal_crud.columns[1]

                    if snr_col_sal_crud:
                        categories_sc = []
                        proceso_vals_sc = []
                        diario_vals_sc = []
                        programa_vals_sc = []
                        columna1_vals_sc = []

                        prod_rows_sc = []
                        for idx, row in app.df_prod_sal_crud.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sc.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sc.append((cat, val))
                                except: pass

                        if len(prod_rows_sc) > 30: prod_rows_sc = prod_rows_sc[-30:]

                        for i in range(len(prod_rows_sc)):
                            categories_sc.append(prod_rows_sc[i][0])
                            try: proceso_vals_sc.append(float(prod_rows_sc[i][1]))
                            except: proceso_vals_sc.append(None)
                            diario_vals_sc.append(None)
                            programa_vals_sc.append(None)
                            columna1_vals_sc.append(None)

                        for i in range(31):
                            categories_sc.append(str(i + 1))
                            proceso_vals_sc.append(None)
                            
                            try: diario_vals_sc.append(float(app.df_data_sal_crud[snr_col_sal_crud].iloc[i]))
                            except: diario_vals_sc.append(None)
                            
                            try: programa_vals_sc.append(float(app.df_snr_sal_crud.iloc[i, 0]))
                            except: programa_vals_sc.append(None)
                            
                            try: columna1_vals_sc.append(float(app.df_snr_sal_crud.iloc[i, 1]))
                            except: columna1_vals_sc.append(None)

                        update_slide_chart(chart_sal_crud, categories_sc, proceso_vals_sc, diario_vals_sc, programa_vals_sc, columna1_vals_sc, wine_color, green_color)

            # --- 25. PROCESAR DIAPOSITIVA DE GASOLINAS SALINA CRUZ (DIAPOSITIVA 37) ---
            if len(prs.slides) > 36 and app.df_data_sal_gas is not None and app.df_snr_sal_gas is not None and app.df_prod_sal_gas is not None:
                slide_sal_gas = prs.slides[36]
                chart_sal_gas = None
                for shape in slide_sal_gas.shapes:
                    if shape.has_chart:
                        chart_sal_gas = shape.chart
                        break
                
                if chart_sal_gas:
                    snr_col_sal_gas = None
                    for col in app.df_data_sal_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sal_gas = col
                            break
                    if not snr_col_sal_gas and len(app.df_data_sal_gas.columns) >= 2:
                        snr_col_sal_gas = app.df_data_sal_gas.columns[1]

                    if snr_col_sal_gas:
                        categories_sg = []
                        proceso_vals_sg = []
                        diario_vals_sg = []
                        programa_vals_sg = []
                        columna1_vals_sg = []

                        prod_rows_sg = []
                        for idx, row in app.df_prod_sal_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sg.append((cat, val))
                                except: pass

                        if len(prod_rows_sg) > 30: prod_rows_sg = prod_rows_sg[-30:]

                        for i in range(len(prod_rows_sg)):
                            categories_sg.append(prod_rows_sg[i][0])
                            try: proceso_vals_sg.append(float(prod_rows_sg[i][1]))
                            except: proceso_vals_sg.append(None)
                            diario_vals_sg.append(None)
                            programa_vals_sg.append(None)
                            columna1_vals_sg.append(None)

                        for i in range(31):
                            categories_sg.append(str(i + 1))
                            proceso_vals_sg.append(None)
                            
                            try: diario_vals_sg.append(float(app.df_data_sal_gas[snr_col_sal_gas].iloc[i]))
                            except: diario_vals_sg.append(None)
                            
                            try: programa_vals_sg.append(float(app.df_snr_sal_gas.iloc[i, 0]))
                            except: programa_vals_sg.append(None)
                            
                            try: columna1_vals_sg.append(float(app.df_snr_sal_gas.iloc[i, 1]))
                            except: columna1_vals_sg.append(None)

                        update_slide_chart(chart_sal_gas, categories_sg, proceso_vals_sg, diario_vals_sg, programa_vals_sg, columna1_vals_sg, wine_color, green_color)

            # --- 26. PROCESAR DIAPOSITIVA DE DIESEL SALINA CRUZ (DIAPOSITIVA 38) ---
            if len(prs.slides) > 37 and app.df_data_sal_die is not None and app.df_snr_sal_die is not None and app.df_prod_sal_die is not None:
                slide_sal_die = prs.slides[37]
                chart_sal_die = None
                for shape in slide_sal_die.shapes:
                    if shape.has_chart:
                        chart_sal_die = shape.chart
                        break
                
                if chart_sal_die:
                    snr_col_sal_die = None
                    for col in app.df_data_sal_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sal_die = col
                            break
                    if not snr_col_sal_die and len(app.df_data_sal_die.columns) >= 2:
                        snr_col_sal_die = app.df_data_sal_die.columns[1]

                    if snr_col_sal_die:
                        categories_sd = []
                        proceso_vals_sd = []
                        diario_vals_sd = []
                        programa_vals_sd = []
                        columna1_vals_sd = []

                        prod_rows_sd = []
                        for idx, row in app.df_prod_sal_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sd.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sd.append((cat, val))
                                except: pass

                        if len(prod_rows_sd) > 30: prod_rows_sd = prod_rows_sd[-30:]

                        for i in range(len(prod_rows_sd)):
                            categories_sd.append(prod_rows_sd[i][0])
                            try: proceso_vals_sd.append(float(prod_rows_sd[i][1]))
                            except: proceso_vals_sd.append(None)
                            diario_vals_sd.append(None)
                            programa_vals_sd.append(None)
                            columna1_vals_sd.append(None)

                        for i in range(31):
                            categories_sd.append(str(i + 1))
                            proceso_vals_sd.append(None)
                            
                            try: diario_vals_sd.append(float(app.df_data_sal_die[snr_col_sal_die].iloc[i]))
                            except: diario_vals_sd.append(None)
                            
                            try: programa_vals_sd.append(float(app.df_snr_sal_die.iloc[i, 0]))
                            except: programa_vals_sd.append(None)
                            
                            try: columna1_vals_sd.append(float(app.df_snr_sal_die.iloc[i, 1]))
                            except: columna1_vals_sd.append(None)

                        update_slide_chart(chart_sal_die, categories_sd, proceso_vals_sd, diario_vals_sd, programa_vals_sd, columna1_vals_sd, wine_color, green_color)

            # --- 27. PROCESAR DIAPOSITIVA DE TURBOSINA SALINA CRUZ (DIAPOSITIVA 39) ---
            if len(prs.slides) > 38 and app.df_data_sal_turb is not None and app.df_snr_sal_turb is not None and app.df_prod_sal_turb is not None:
                slide_sal_turb = prs.slides[38]
                chart_sal_turb = None
                for shape in slide_sal_turb.shapes:
                    if shape.has_chart:
                        chart_sal_turb = shape.chart
                        break
                
                if chart_sal_turb:
                    snr_col_sal_turb = None
                    for col in app.df_data_sal_turb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sal_turb = col
                            break
                    if not snr_col_sal_turb and len(app.df_data_sal_turb.columns) >= 2:
                        snr_col_sal_turb = app.df_data_sal_turb.columns[1]

                    if snr_col_sal_turb:
                        categories_stur = []
                        proceso_vals_stur = []
                        diario_vals_stur = []
                        programa_vals_stur = []
                        columna1_vals_stur = []

                        prod_rows_stur = []
                        for idx, row in app.df_prod_sal_turb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_stur.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_stur.append((cat, val))
                                except: pass

                        if len(prod_rows_stur) > 30: prod_rows_stur = prod_rows_stur[-30:]

                        for i in range(len(prod_rows_stur)):
                            categories_stur.append(prod_rows_stur[i][0])
                            try: proceso_vals_stur.append(float(prod_rows_stur[i][1]))
                            except: proceso_vals_stur.append(None)
                            diario_vals_stur.append(None)
                            programa_vals_stur.append(None)
                            columna1_vals_stur.append(None)

                        for i in range(31):
                            categories_stur.append(str(i + 1))
                            proceso_vals_stur.append(None)
                            
                            try: diario_vals_stur.append(float(app.df_data_sal_turb[snr_col_sal_turb].iloc[i]))
                            except: diario_vals_stur.append(None)
                            
                            try: programa_vals_stur.append(float(app.df_snr_sal_turb.iloc[i, 0]))
                            except: programa_vals_stur.append(None)
                            
                            try: columna1_vals_stur.append(float(app.df_snr_sal_turb.iloc[i, 1]))
                            except: columna1_vals_stur.append(None)

                        update_slide_chart(chart_sal_turb, categories_stur, proceso_vals_stur, diario_vals_stur, programa_vals_stur, columna1_vals_stur, wine_color, green_color)

            # --- 28. PROCESAR DIAPOSITIVA DE COMBUSTOLEO SALINA CRUZ (DIAPOSITIVA 40) ---
            if len(prs.slides) > 39 and app.df_data_sal_comb is not None and app.df_snr_sal_comb is not None and app.df_prod_sal_comb is not None:
                slide_sal_comb = prs.slides[39]
                chart_sal_comb = None
                for shape in slide_sal_comb.shapes:
                    if shape.has_chart:
                        chart_sal_comb = shape.chart
                        break
                
                if chart_sal_comb:
                    snr_col_sal_comb = None
                    for col in app.df_data_sal_comb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_sal_comb = col
                            break
                    if not snr_col_sal_comb and len(app.df_data_sal_comb.columns) >= 2:
                        snr_col_sal_comb = app.df_data_sal_comb.columns[1]

                    if snr_col_sal_comb:
                        categories_sco = []
                        proceso_vals_sco = []
                        diario_vals_sco = []
                        programa_vals_sco = []
                        columna1_vals_sco = []

                        prod_rows_sco = []
                        for idx, row in app.df_prod_sal_comb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sco.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sco.append((cat, val))
                                except: pass

                        if len(prod_rows_sco) > 30: prod_rows_sco = prod_rows_sco[-30:]

                        for i in range(len(prod_rows_sco)):
                            categories_sco.append(prod_rows_sco[i][0])
                            try: proceso_vals_sco.append(float(prod_rows_sco[i][1]))
                            except: proceso_vals_sco.append(None)
                            diario_vals_sco.append(None)
                            programa_vals_sco.append(None)
                            columna1_vals_sco.append(None)

                        for i in range(31):
                            categories_sco.append(str(i + 1))
                            proceso_vals_sco.append(None)
                            
                            try: diario_vals_sco.append(float(app.df_data_sal_comb[snr_col_sal_comb].iloc[i]))
                            except: diario_vals_sco.append(None)
                            
                            try: programa_vals_sco.append(float(app.df_snr_sal_comb.iloc[i, 0]))
                            except: programa_vals_sco.append(None)
                            
                            try: columna1_vals_sco.append(float(app.df_snr_sal_comb.iloc[i, 1]))
                            except: columna1_vals_sco.append(None)

                        update_slide_chart(chart_sal_comb, categories_sco, proceso_vals_sco, diario_vals_sco, programa_vals_sco, columna1_vals_sco, wine_color, green_color)

            # --- 29. PROCESAR DIAPOSITIVA DE CRUDO TULA (DIAPOSITIVA 43) ---
            if len(prs.slides) > 42 and app.df_data_tula_crud is not None and app.df_snr_tula_crud is not None and app.df_prod_tula_crud is not None:
                slide_tula_crud = prs.slides[42]
                chart_tula_crud = None
                for shape in slide_tula_crud.shapes:
                    if shape.has_chart:
                        chart_tula_crud = shape.chart
                        break
                
                if chart_tula_crud:
                    snr_col_tula_crud = None
                    for col in app.df_data_tula_crud.columns:
                        if "SNR" in str(col).upper():
                            snr_col_tula_crud = col
                            break
                    if not snr_col_tula_crud and len(app.df_data_tula_crud.columns) >= 2:
                        snr_col_tula_crud = app.df_data_tula_crud.columns[1]

                    if snr_col_tula_crud:
                        categories_sc = []
                        proceso_vals_sc = []
                        diario_vals_sc = []
                        programa_vals_sc = []
                        columna1_vals_sc = []

                        prod_rows_sc = []
                        for idx, row in app.df_prod_tula_crud.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sc.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sc.append((cat, val))
                                except: pass

                        if len(prod_rows_sc) > 30: prod_rows_sc = prod_rows_sc[-30:]

                        for i in range(len(prod_rows_sc)):
                            categories_sc.append(prod_rows_sc[i][0])
                            try: proceso_vals_sc.append(float(prod_rows_sc[i][1]))
                            except: proceso_vals_sc.append(None)
                            diario_vals_sc.append(None)
                            programa_vals_sc.append(None)
                            columna1_vals_sc.append(None)

                        for i in range(31):
                            categories_sc.append(str(i + 1))
                            proceso_vals_sc.append(None)
                            
                            try: diario_vals_sc.append(float(app.df_data_tula_crud[snr_col_tula_crud].iloc[i]))
                            except: diario_vals_sc.append(None)
                            
                            try: programa_vals_sc.append(float(app.df_snr_tula_crud.iloc[i, 0]))
                            except: programa_vals_sc.append(None)
                            
                            try: columna1_vals_sc.append(float(app.df_snr_tula_crud.iloc[i, 1]))
                            except: columna1_vals_sc.append(None)

                        update_slide_chart(chart_tula_crud, categories_sc, proceso_vals_sc, diario_vals_sc, programa_vals_sc, columna1_vals_sc, wine_color, green_color)

            # --- 30. PROCESAR DIAPOSITIVA DE GASOLINAS TULA (DIAPOSITIVA 44) ---
            if len(prs.slides) > 43 and app.df_data_tula_gas is not None and app.df_snr_tula_gas is not None and app.df_prod_tula_gas is not None:
                slide_tula_gas = prs.slides[43]
                chart_tula_gas = None
                for shape in slide_tula_gas.shapes:
                    if shape.has_chart:
                        chart_tula_gas = shape.chart
                        break
                
                if chart_tula_gas:
                    snr_col_tula_gas = None
                    for col in app.df_data_tula_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_tula_gas = col
                            break
                    if not snr_col_tula_gas and len(app.df_data_tula_gas.columns) >= 2:
                        snr_col_tula_gas = app.df_data_tula_gas.columns[1]

                    if snr_col_tula_gas:
                        categories_sg = []
                        proceso_vals_sg = []
                        diario_vals_sg = []
                        programa_vals_sg = []
                        columna1_vals_sg = []

                        prod_rows_sg = []
                        for idx, row in app.df_prod_tula_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sg.append((cat, val))
                                except: pass

                        if len(prod_rows_sg) > 30: prod_rows_sg = prod_rows_sg[-30:]

                        for i in range(len(prod_rows_sg)):
                            categories_sg.append(prod_rows_sg[i][0])
                            try: proceso_vals_sg.append(float(prod_rows_sg[i][1]))
                            except: proceso_vals_sg.append(None)
                            diario_vals_sg.append(None)
                            programa_vals_sg.append(None)
                            columna1_vals_sg.append(None)

                        for i in range(31):
                            categories_sg.append(str(i + 1))
                            proceso_vals_sg.append(None)
                            
                            try: diario_vals_sg.append(float(app.df_data_tula_gas[snr_col_tula_gas].iloc[i]))
                            except: diario_vals_sg.append(None)
                            
                            try: programa_vals_sg.append(float(app.df_snr_tula_gas.iloc[i, 0]))
                            except: programa_vals_sg.append(None)
                            
                            try: columna1_vals_sg.append(float(app.df_snr_tula_gas.iloc[i, 1]))
                            except: columna1_vals_sg.append(None)

                        update_slide_chart(chart_tula_gas, categories_sg, proceso_vals_sg, diario_vals_sg, programa_vals_sg, columna1_vals_sg, wine_color, green_color)

            # --- 31. PROCESAR DIAPOSITIVA DE DIESEL TULA (DIAPOSITIVA 45) ---
            if len(prs.slides) > 44 and app.df_data_tula_die is not None and app.df_snr_tula_die is not None and app.df_prod_tula_die is not None:
                slide_tula_die = prs.slides[44]
                chart_tula_die = None
                for shape in slide_tula_die.shapes:
                    if shape.has_chart:
                        chart_tula_die = shape.chart
                        break
                
                if chart_tula_die:
                    snr_col_tula_die = None
                    for col in app.df_data_tula_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_tula_die = col
                            break
                    if not snr_col_tula_die and len(app.df_data_tula_die.columns) >= 2:
                        snr_col_tula_die = app.df_data_tula_die.columns[1]

                    if snr_col_tula_die:
                        categories_sd = []
                        proceso_vals_sd = []
                        diario_vals_sd = []
                        programa_vals_sd = []
                        columna1_vals_sd = []

                        prod_rows_sd = []
                        for idx, row in app.df_prod_tula_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sd.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sd.append((cat, val))
                                except: pass

                        if len(prod_rows_sd) > 30: prod_rows_sd = prod_rows_sd[-30:]

                        for i in range(len(prod_rows_sd)):
                            categories_sd.append(prod_rows_sd[i][0])
                            try: proceso_vals_sd.append(float(prod_rows_sd[i][1]))
                            except: proceso_vals_sd.append(None)
                            diario_vals_sd.append(None)
                            programa_vals_sd.append(None)
                            columna1_vals_sd.append(None)

                        for i in range(31):
                            categories_sd.append(str(i + 1))
                            proceso_vals_sd.append(None)
                            
                            try: diario_vals_sd.append(float(app.df_data_tula_die[snr_col_tula_die].iloc[i]))
                            except: diario_vals_sd.append(None)
                            
                            try: programa_vals_sd.append(float(app.df_snr_tula_die.iloc[i, 0]))
                            except: programa_vals_sd.append(None)
                            
                            try: columna1_vals_sd.append(float(app.df_snr_tula_die.iloc[i, 1]))
                            except: columna1_vals_sd.append(None)

                        update_slide_chart(chart_tula_die, categories_sd, proceso_vals_sd, diario_vals_sd, programa_vals_sd, columna1_vals_sd, wine_color, green_color)

            # --- 32. PROCESAR DIAPOSITIVA DE TURBOSINA TULA (DIAPOSITIVA 46) ---
            if len(prs.slides) > 45 and app.df_data_tula_turb is not None and app.df_snr_tula_turb is not None and app.df_prod_tula_turb is not None:
                slide_tula_turb = prs.slides[45]
                chart_tula_turb = None
                for shape in slide_tula_turb.shapes:
                    if shape.has_chart:
                        chart_tula_turb = shape.chart
                        break
                
                if chart_tula_turb:
                    snr_col_tula_turb = None
                    for col in app.df_data_tula_turb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_tula_turb = col
                            break
                    if not snr_col_tula_turb and len(app.df_data_tula_turb.columns) >= 2:
                        snr_col_tula_turb = app.df_data_tula_turb.columns[1]

                    if snr_col_tula_turb:
                        categories_stur = []
                        proceso_vals_stur = []
                        diario_vals_stur = []
                        programa_vals_stur = []
                        columna1_vals_stur = []

                        prod_rows_stur = []
                        for idx, row in app.df_prod_tula_turb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_stur.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_stur.append((cat, val))
                                except: pass

                        if len(prod_rows_stur) > 30: prod_rows_stur = prod_rows_stur[-30:]

                        for i in range(len(prod_rows_stur)):
                            categories_stur.append(prod_rows_stur[i][0])
                            try: proceso_vals_stur.append(float(prod_rows_stur[i][1]))
                            except: proceso_vals_stur.append(None)
                            diario_vals_stur.append(None)
                            programa_vals_stur.append(None)
                            columna1_vals_stur.append(None)

                        for i in range(31):
                            categories_stur.append(str(i + 1))
                            proceso_vals_stur.append(None)
                            
                            try: diario_vals_stur.append(float(app.df_data_tula_turb[snr_col_tula_turb].iloc[i]))
                            except: diario_vals_stur.append(None)
                            
                            try: programa_vals_stur.append(float(app.df_snr_tula_turb.iloc[i, 0]))
                            except: programa_vals_stur.append(None)
                            
                            try: columna1_vals_stur.append(float(app.df_snr_tula_turb.iloc[i, 1]))
                            except: columna1_vals_stur.append(None)

                        update_slide_chart(chart_tula_turb, categories_stur, proceso_vals_stur, diario_vals_stur, programa_vals_stur, columna1_vals_stur, wine_color, green_color)

            # --- 33. PROCESAR DIAPOSITIVA DE COMBUSTOLEO TULA (DIAPOSITIVA 47) ---
            if len(prs.slides) > 46 and app.df_data_tula_comb is not None and app.df_snr_tula_comb is not None and app.df_prod_tula_comb is not None:
                slide_tula_comb = prs.slides[46]
                chart_tula_comb = None
                for shape in slide_tula_comb.shapes:
                    if shape.has_chart:
                        chart_tula_comb = shape.chart
                        break
                
                if chart_tula_comb:
                    snr_col_tula_comb = None
                    for col in app.df_data_tula_comb.columns:
                        if "SNR" in str(col).upper():
                            snr_col_tula_comb = col
                            break
                    if not snr_col_tula_comb and len(app.df_data_tula_comb.columns) >= 2:
                        snr_col_tula_comb = app.df_data_tula_comb.columns[1]

                    if snr_col_tula_comb:
                        categories_sco = []
                        proceso_vals_sco = []
                        diario_vals_sco = []
                        programa_vals_sco = []
                        columna1_vals_sco = []

                        prod_rows_sco = []
                        for idx, row in app.df_prod_tula_comb.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sco.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sco.append((cat, val))
                                except: pass

                        if len(prod_rows_sco) > 30: prod_rows_sco = prod_rows_sco[-30:]

                        for i in range(len(prod_rows_sco)):
                            categories_sco.append(prod_rows_sco[i][0])
                            try: proceso_vals_sco.append(float(prod_rows_sco[i][1]))
                            except: proceso_vals_sco.append(None)
                            diario_vals_sco.append(None)
                            programa_vals_sco.append(None)
                            columna1_vals_sco.append(None)

                        for i in range(31):
                            categories_sco.append(str(i + 1))
                            proceso_vals_sco.append(None)
                            
                            try: diario_vals_sco.append(float(app.df_data_tula_comb[snr_col_tula_comb].iloc[i]))
                            except: diario_vals_sco.append(None)
                            
                            try: programa_vals_sco.append(float(app.df_snr_tula_comb.iloc[i, 0]))
                            except: programa_vals_sco.append(None)
                            
                            try: columna1_vals_sco.append(float(app.df_snr_tula_comb.iloc[i, 1]))
                            except: columna1_vals_sco.append(None)

                        update_slide_chart(chart_tula_comb, categories_sco, proceso_vals_sco, diario_vals_sco, programa_vals_sco, columna1_vals_sco, wine_color, green_color)

            # --- 34. PROCESAR DIAPOSITIVA DE CRUDO OLMECA (DIAPOSITIVA 50) ---
            if len(prs.slides) > 49 and app.df_data_olme_crud is not None and app.df_snr_olme_crud is not None and app.df_prod_olme_crud is not None:
                slide_olme_crud = prs.slides[49]
                chart_olme_crud = None
                for shape in slide_olme_crud.shapes:
                    if shape.has_chart:
                        chart_olme_crud = shape.chart
                        break
                
                if chart_olme_crud:
                    snr_col_olme_crud = None
                    for col in app.df_data_olme_crud.columns:
                        if "SNR" in str(col).upper():
                            snr_col_olme_crud = col
                            break
                    if not snr_col_olme_crud and len(app.df_data_olme_crud.columns) >= 2:
                        snr_col_olme_crud = app.df_data_olme_crud.columns[1]

                    if snr_col_olme_crud:
                        categories_sc = []
                        proceso_vals_sc = []
                        diario_vals_sc = []
                        programa_vals_sc = []
                        columna1_vals_sc = []

                        prod_rows_sc = []
                        for idx, row in app.df_prod_olme_crud.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sc.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sc.append((cat, val))
                                except: pass

                        if len(prod_rows_sc) > 30: prod_rows_sc = prod_rows_sc[-30:]

                        for i in range(len(prod_rows_sc)):
                            categories_sc.append(prod_rows_sc[i][0])
                            try: proceso_vals_sc.append(float(prod_rows_sc[i][1]))
                            except: proceso_vals_sc.append(None)
                            diario_vals_sc.append(None)
                            programa_vals_sc.append(None)
                            columna1_vals_sc.append(None)

                        for i in range(31):
                            categories_sc.append(str(i + 1))
                            proceso_vals_sc.append(None)
                            
                            try: diario_vals_sc.append(float(app.df_data_olme_crud[snr_col_olme_crud].iloc[i]))
                            except: diario_vals_sc.append(None)
                            
                            try: programa_vals_sc.append(float(app.df_snr_olme_crud.iloc[i, 0]))
                            except: programa_vals_sc.append(None)
                            
                            try: columna1_vals_sc.append(float(app.df_snr_olme_crud.iloc[i, 1]))
                            except: columna1_vals_sc.append(None)

                        update_slide_chart(chart_olme_crud, categories_sc, proceso_vals_sc, diario_vals_sc, programa_vals_sc, columna1_vals_sc, wine_color, green_color)

            # --- 35. PROCESAR DIAPOSITIVA DE GASOLINAS OLMECA (DIAPOSITIVA 51) ---
            if len(prs.slides) > 50 and app.df_data_olme_gas is not None and app.df_snr_olme_gas is not None and app.df_prod_olme_gas is not None:
                slide_olme_gas = prs.slides[50]
                chart_olme_gas = None
                for shape in slide_olme_gas.shapes:
                    if shape.has_chart:
                        chart_olme_gas = shape.chart
                        break
                
                if chart_olme_gas:
                    snr_col_olme_gas = None
                    for col in app.df_data_olme_gas.columns:
                        if "SNR" in str(col).upper():
                            snr_col_olme_gas = col
                            break
                    if not snr_col_olme_gas and len(app.df_data_olme_gas.columns) >= 2:
                        snr_col_olme_gas = app.df_data_olme_gas.columns[1]

                    if snr_col_olme_gas:
                        categories_sg = []
                        proceso_vals_sg = []
                        diario_vals_sg = []
                        programa_vals_sg = []
                        columna1_vals_sg = []

                        prod_rows_sg = []
                        for idx, row in app.df_prod_olme_gas.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sg.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sg.append((cat, val))
                                except: pass

                        if len(prod_rows_sg) > 30: prod_rows_sg = prod_rows_sg[-30:]

                        for i in range(len(prod_rows_sg)):
                            categories_sg.append(prod_rows_sg[i][0])
                            try: proceso_vals_sg.append(float(prod_rows_sg[i][1]))
                            except: proceso_vals_sg.append(None)
                            diario_vals_sg.append(None)
                            programa_vals_sg.append(None)
                            columna1_vals_sg.append(None)

                        for i in range(31):
                            categories_sg.append(str(i + 1))
                            proceso_vals_sg.append(None)
                            
                            try: diario_vals_sg.append(float(app.df_data_olme_gas[snr_col_olme_gas].iloc[i]))
                            except: diario_vals_sg.append(None)
                            
                            try: programa_vals_sg.append(float(app.df_snr_olme_gas.iloc[i, 0]))
                            except: programa_vals_sg.append(None)
                            
                            try: columna1_vals_sg.append(float(app.df_snr_olme_gas.iloc[i, 1]))
                            except: columna1_vals_sg.append(None)

                        update_slide_chart(chart_olme_gas, categories_sg, proceso_vals_sg, diario_vals_sg, programa_vals_sg, columna1_vals_sg, wine_color, green_color)

            # --- 36. PROCESAR DIAPOSITIVA DE DIESEL OLMECA (DIAPOSITIVA 52) ---
            if len(prs.slides) > 51 and app.df_data_olme_die is not None and app.df_snr_olme_die is not None and app.df_prod_olme_die is not None:
                slide_olme_die = prs.slides[51]
                chart_olme_die = None
                for shape in slide_olme_die.shapes:
                    if shape.has_chart:
                        chart_olme_die = shape.chart
                        break
                
                if chart_olme_die:
                    snr_col_olme_die = None
                    for col in app.df_data_olme_die.columns:
                        if "SNR" in str(col).upper():
                            snr_col_olme_die = col
                            break
                    if not snr_col_olme_die and len(app.df_data_olme_die.columns) >= 2:
                        snr_col_olme_die = app.df_data_olme_die.columns[1]

                    if snr_col_olme_die:
                        categories_sd = []
                        proceso_vals_sd = []
                        diario_vals_sd = []
                        programa_vals_sd = []
                        columna1_vals_sd = []

                        prod_rows_sd = []
                        for idx, row in app.df_prod_olme_die.iterrows():
                            cat = str(row.iloc[0]).strip()
                            val = row.iloc[1]
                            if not cat: continue
                            if not any(c.isalpha() for c in cat):
                                prod_rows_sd.append((cat, val))
                            else:
                                try:
                                    if float(val) != 0: prod_rows_sd.append((cat, val))
                                except: pass

                        if len(prod_rows_sd) > 30: prod_rows_sd = prod_rows_sd[-30:]

                        for i in range(len(prod_rows_sd)):
                            categories_sd.append(prod_rows_sd[i][0])
                            try: proceso_vals_sd.append(float(prod_rows_sd[i][1]))
                            except: proceso_vals_sd.append(None)
                            diario_vals_sd.append(None)
                            programa_vals_sd.append(None)
                            columna1_vals_sd.append(None)

                        for i in range(31):
                            categories_sd.append(str(i + 1))
                            proceso_vals_sd.append(None)
                            
                            try: diario_vals_sd.append(float(app.df_data_olme_die[snr_col_olme_die].iloc[i]))
                            except: diario_vals_sd.append(None)
                            
                            try: programa_vals_sd.append(float(app.df_snr_olme_die.iloc[i, 0]))
                            except: programa_vals_sd.append(None)
                            
                            try: columna1_vals_sd.append(float(app.df_snr_olme_die.iloc[i, 1]))
                            except: columna1_vals_sd.append(None)

                        update_slide_chart(chart_olme_die, categories_sd, proceso_vals_sd, diario_vals_sd, programa_vals_sd, columna1_vals_sd, wine_color, green_color)

        prs.save(save_path)


        app.after(0, app.on_pptx_success, save_path)

    except Exception as e:
        err_details = traceback.format_exc()
        app.after(0, app.on_pptx_error, str(e), err_details)

