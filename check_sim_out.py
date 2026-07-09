from pptx import Presentation

prs = Presentation("simulate_out.pptx")
slide = prs.slides[1]
chart = slide.shapes[0].chart
series = chart.series[0]
for p_idx in range(17):
    point = series.points[p_idx]
    fill = point.format.fill
    print(f"Point {p_idx} ({chart.plots[0].categories[p_idx].label}): fill.type={fill.type}")
    try:
        print(f"  rgb={fill.fore_color.rgb}")
    except Exception as e:
        print(f"  rgb error: {e}")
